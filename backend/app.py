import os
import hashlib
import requests # Used for OpenLibrary API calls
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, timedelta
import random
from urllib.parse import urljoin
from flask import Flask, request, jsonify, Blueprint
from flask_sqlalchemy import SQLAlchemy
from flask_cors import CORS
from flask_jwt_extended import (
    JWTManager, create_access_token, jwt_required, get_jwt_identity, get_jwt,
    verify_jwt_in_request, get_jwt_identity as _get_jwt_identity
)
from werkzeug.utils import secure_filename
import json
import time
import uuid
import io
import csv
import boto3
from functools import wraps
# File Processing Libs
try:
    import docx
except ImportError:
    docx = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
    
# Optional PDF tool
try:
    import pdfkit
except ImportError:
    pdfkit = None
    print("Warning: pdfkit not installed. Using ReportLab/TXT fallback for receipts.")
try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
except ImportError:
    canvas = None
    A4 = None
    print("Warning: reportlab not installed. Receipts will be TXT only.")
from sqlalchemy import func, or_
from dotenv import load_dotenv

# Load .env
load_dotenv()

# -------------------- CONFIG (from .env) --------------------
JWT_SECRET_KEY = os.getenv("JWT_SECRET_KEY", "dev_secret_jwt_key_please_set")
SECRET_KEY = os.getenv("SECRET_KEY") or JWT_SECRET_KEY
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
PASSWORD_SALT = os.getenv("PASSWORD_SALT", "noteorbit_salt_v1")
DEFAULT_ADMIN_EMAIL = os.getenv("DEFAULT_ADMIN_EMAIL", "admin@noteorbit.edu")
DEFAULT_ADMIN_PASSWORD = os.getenv("DEFAULT_ADMIN_PASSWORD", "admin123")
DEFAULT_ADMIN_PASSWORD = os.getenv("DEFAULT_ADMIN_PASSWORD", "admin123")
HRD_EMAIL = os.getenv("HRD_EMAIL", "hrd@noteorbit.edu")
HRD_PASSWORD = os.getenv("HRD_PASSWORD", "hrdsnpsu123")
FLASK_RUN_PORT = int(os.getenv("FLASK_RUN_PORT", 5000))
SMTP_HOST = os.getenv("SMTP_HOST")
SMTP_PORT = int(os.getenv("SMTP_PORT", 587))
SMTP_USER = os.getenv("SMTP_USER")
SMTP_PASS = os.getenv("SMTP_PASS")

# -------------------- HARDCODED MINIO CREDENTIALS (kept in-code per request) --------------------
S3_ENDPOINT = "https://bulk-automatic-groove-sage.trycloudflare.com"
S3_BUCKET = "noteorbit"
S3_ACCESS_KEY = "admin"
S3_SECRET_KEY = "password123"
# -------------------- END MINIO CREDENTIALS --------------------

# -------------------- EXTERNAL API CONFIG --------------------
OPENLIBRARY_URL = "https://openlibrary.org/search.json"

# -------------------- APP INIT --------------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMP_DIR = os.path.join(BASE_DIR, "tmp")
RECEIPT_TMP_DIR = os.path.join(TEMP_DIR, "receipts")
os.makedirs(RECEIPT_TMP_DIR, exist_ok=True)

app = Flask(__name__)
app.config["SECRET_KEY"] = SECRET_KEY
app.config["JWT_SECRET_KEY"] = JWT_SECRET_KEY
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024
# CORS - allow local frontend origins (adjust if needed)
CORS(
    app,
    resources={r"*": {
        "origins": [
            "https://note-orbit.vercel.app",
            "*",
        ],
        "allow_headers": ["Content-Type", "Authorization"],
        "supports_credentials": True,
        "methods": ["GET", "POST", "PUT", "PATCH", "DELETE", "OPTIONS"]
    }}
)

# DB
DB_PATH = os.path.join(BASE_DIR, "noteorbit.db")
app.config["SQLALCHEMY_DATABASE_URI"] = f"sqlite:///{DB_PATH}"
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False

db = SQLAlchemy(app)
jwt = JWTManager(app)

SALT = PASSWORD_SALT
RECEIPT_PREFIX = "fees/"

# -------------------- S3 / MinIO client --------------------
s3_client = boto3.client(
    "s3",
    endpoint_url=S3_ENDPOINT,
    aws_access_key_id=S3_ACCESS_KEY,
    aws_secret_access_key=S3_SECRET_KEY,
    region_name="us-east-1",
)

# -------------------- HELPERS --------------------

def send_email(to_email, subject, body):
    if not SMTP_HOST or not SMTP_USER or not SMTP_PASS:
        print("SMTP not configured. Skipping email.")
        return False
    try:
        msg = MIMEMultipart()
        msg['From'] = SMTP_USER
        msg['To'] = to_email
        msg['Subject'] = subject
        msg.attach(MIMEText(body, 'html')) # Changed to HTML for better formatting

        server = smtplib.SMTP(SMTP_HOST, SMTP_PORT)
        server.starttls()
        server.login(SMTP_USER, SMTP_PASS)
        server.sendmail(SMTP_USER, to_email, msg.as_string())
        server.quit()
        return True
    except Exception as e:
        print(f"Failed to send email: {e}")
        return False

def send_professional_email(to_email, subject, title, details, main_body):
    """
    Sends a professionally styled HTML email.
    details: dict of {'Label': 'Value'}
    """
    if not to_email:
        return False
        
    html_content = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <style>
            body {{ font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif; background-color: #f4f4f9; color: #333; margin: 0; padding: 0; }}
            .container {{ max-width: 600px; margin: 20px auto; background: #ffffff; border-radius: 12px; box-shadow: 0 4px 15px rgba(0,0,0,0.1); overflow: hidden; border-left: 5px solid #4f46e5; }}
            .header {{ background-color: #fff; padding: 30px 40px; border-bottom: 1px solid #eee; }}
            .header h1 {{ margin: 0; color: #4f46e5; font-size: 24px; letter-spacing: -0.5px; }}
            .header p {{ margin: 5px 0 0; color: #6b7280; font-size: 14px; text-transform: uppercase; letter-spacing: 1px; font-weight: 600; }}
            .content {{ padding: 30px 40px; line-height: 1.6; color: #374151; }}
            .details-box {{ background-color: #f9fafb; border: 1px solid #e5e7eb; border-radius: 8px; padding: 20px; margin: 25px 0; }}
            .detail-row {{ display: flex; justify-content: space-between; padding: 8px 0; border-bottom: 1px solid #eee; }}
            .detail-row:last-child {{ border-bottom: none; }}
            .label {{ color: #6b7280; font-weight: 500; font-size: 13px; }}
            .value {{ color: #111827; font-weight: 600; font-size: 14px; text-align: right; }}
            .footer {{ background-color: #f9fafb; padding: 20px 40px; text-align: center; font-size: 12px; color: #9ca3af; border-top: 1px solid #eee; }}
            .btn {{ display: inline-block; background-color: #4f46e5; color: white; padding: 10px 20px; text-decoration: none; border-radius: 6px; font-size: 14px; font-weight: 600; margin-top: 10px; }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">
                <p>NoteOrbit Notification</p>
                <h1>{title}</h1>
            </div>
            <div class="content">
                <p>Dear User,</p>
                <p>{main_body}</p>
                
                <div class="details-box">
                    {"".join([f'<div class="detail-row"><span class="label">{k}</span><span class="value">{v}</span></div>' for k, v in details.items()])}
                </div>
                
                <p>Please log in to the portal to view full details.</p>
                <center><a href="http://localhost:5173" class="btn">Login to NoteOrbit</a></center>
            </div>
            <div class="footer">
                &copy; 2025 NoteOrbit Academic System by LeafCore Labs. All rights reserved.<br>
                This is an automated message. Please do not reply.
            </div>
        </div>
    </body>
    </html>
    """
    return send_email(to_email, subject, html_content)




def generate_otp():
    return str(random.randint(100000, 999999))


def hash_password(password: str, salt: str = SALT) -> str:
    return hashlib.sha256((password + salt).encode()).hexdigest()


def allowed_file(filename: str):
    ALLOWED = {"pdf", "doc", "docx", "ppt", "pptx", "txt", "epub", "jpg", "jpeg", "png", "csv"}
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED


def upload_to_minio(file_obj, dest_key: str, content_type: str = "application/octet-stream"):
    """Handles file upload to MinIO (S3 compatible). Returns presigned URL and key."""
    try:
        # Ensure bucket exists
        s3_client.head_bucket(Bucket=S3_BUCKET)
    except Exception:
        try:
            s3_client.create_bucket(Bucket=S3_BUCKET)
        except Exception:
            pass

    # If file_obj is a werkzeug FileStorage, it has .stream; but upload_fileobj accepts file-like
    s3_client.upload_fileobj(
        file_obj,
        S3_BUCKET,
        dest_key,
        ExtraArgs={"ContentType": content_type},
    )

    url = s3_client.generate_presigned_url(
        "get_object",
        Params={"Bucket": S3_BUCKET, "Key": dest_key},
        ExpiresIn=3600,
    )
    return url, dest_key


def cents_to_rupees_str(amount_cents):
    return f"{amount_cents // 100}.{amount_cents % 100:02d}"


def ensure_receipt_pdf(html_content, filename_base):
    fname = f"{filename_base}.pdf"
    local_path = os.path.join(RECEIPT_TMP_DIR, fname)

    # Try using pdfkit first
    if pdfkit:
        try:
            pdf_bytes = pdfkit.from_string(html_content, False)
            with open(local_path, "wb") as f:
                f.write(pdf_bytes)
            return local_path
        except Exception:
            pass

    # Fallback to ReportLab (Simple PDF)
    try:
        c = canvas.Canvas(local_path, pagesize=A4)
        textobject = c.beginText(40, 800)

        # Simple parsing to extract plain text lines from demo HTML
        temp_content = html_content
        if "<body>" in html_content:
            try:
                temp_content = html_content.split("<body>")[1].split("</body>")[0]
            except Exception:
                temp_content = html_content
        lines = [line.split(">", 1)[-1].replace("</p>", "").replace("</strong>", "") for line in temp_content.split("<p>")]

        for line in lines:
            line_plain = line.strip().replace("<h2>", "").replace("</h2>", "").replace("<strong>", "").replace("</small>", "").replace("<a>", "").replace("</a>", "")
            if line_plain:
                textobject.textLine(line_plain[:200])

        c.drawText(textobject)
        c.showPage()
        c.save()
        return local_path
    except Exception:
        # Final fallback to TXT
        txt_path = os.path.join(RECEIPT_TMP_DIR, f"{filename_base}.txt")
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write("Payment Receipt\n\n")
            import re
            clean = re.sub('<[^>]+>', '\n', html_content)
            f.write(clean)
        return txt_path


def upload_receipt(path_local, dest_key=None):
    if not dest_key:
        dest_key = RECEIPT_PREFIX + os.path.basename(path_local)
    # Determine content type based on extension
    content_type = "application/pdf" if path_local.endswith(".pdf") else "text/plain"

    with open(path_local, "rb") as f:
        upload_to_minio(f, dest_key, content_type)

    presigned = s3_client.generate_presigned_url(
        "get_object",
        Params={"Bucket": S3_BUCKET, "Key": dest_key},
        ExpiresIn=3600,
    )
    return dest_key, presigned


# -------------------- LIBRARY HELPER --------------------

def search_open_library(query: str):
    """Fetches book data from OpenLibrary API and standardizes the output."""
    try:
        # Fetch up to 10 results from OpenLibrary
        response = requests.get(OPENLIBRARY_URL, params={"q": query, "limit": 10})
        response.raise_for_status()
        data = response.json()
        books = []
        for doc in data.get("docs", []):
            # Standardize the output format for the frontend
            author_names = doc.get("author_name")
            
            if doc.get("title") and author_names:
                # Use the first ISBN as a unique identifier if available
                isbn = doc.get("isbn")
                
                books.append({
                    "id": f"OL-{doc.get('key')}",
                    "title": doc["title"],
                    "author": ", ".join(author_names),
                    "source": "OpenLibrary",
                    # Generate a cover URL if cover ID is present
                    "cover_url": f"https://covers.openlibrary.org/b/id/{doc.get('cover_i')}-M.jpg" if doc.get('cover_i') else None,
                    "isbn": isbn[0] if isbn and isinstance(isbn, list) else None,
                    "file_url": None, # External source has no direct download link
                    "degree": None,
                    "semester": None
                })
        return books
    except Exception as e:
        print(f"Error calling OpenLibrary API: {e}")
        return []


# -------------------- DB MODELS --------------------

# --- NEW HOSTEL DB MODELS ---
class Hostel(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), unique=True, nullable=False)
    address = db.Column(db.String(255))
    vacant_beds = db.Column(db.Integer, default=0)

class OTP(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(120), nullable=False)
    otp = db.Column(db.String(6), nullable=False)
    expires_at = db.Column(db.DateTime, nullable=False)
    verified = db.Column(db.Boolean, default=False)

class Room(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    hostel_id = db.Column(db.Integer, db.ForeignKey('hostel.id', ondelete='CASCADE'))
    room_number = db.Column(db.String(20), nullable=False)
    capacity = db.Column(db.Integer, default=1)
    current_occupancy = db.Column(db.Integer, default=0)
    __table_args__ = (
        db.UniqueConstraint('hostel_id', 'room_number', name='uq_room_hostel'),
    )

class HostelAllocation(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey('user.id', ondelete='CASCADE'), unique=True, nullable=False)
    hostel_id = db.Column(db.Integer, db.ForeignKey('hostel.id'))
    room_id = db.Column(db.Integer, db.ForeignKey('room.id'))
    allocated_on = db.Column(db.DateTime, default=datetime.utcnow)

class AIChatSession(db.Model):
    id = db.Column(db.String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(200), default="New Chat")
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow)
    messages = db.relationship('AIChatMessage', backref='session', cascade='all, delete-orphan')

class AIChatMessage(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    session_id = db.Column(db.String(36), db.ForeignKey('ai_chat_session.id', ondelete='CASCADE'), nullable=False)
    role = db.Column(db.String(20), nullable=False) # 'user' or 'ai'
    text = db.Column(db.Text, nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)
    # Optional: store file reference if we want to retrieve attachments later
    attachment = db.Column(db.String(255), nullable=True)

# --- END NEW HOSTEL DB MODELS ---


class Degree(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(50), unique=True, nullable=False)


class Section(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    degree = db.Column(db.String(50), nullable=False)
    semester = db.Column(db.Integer, nullable=False)
    name = db.Column(db.String(20), nullable=False)
    __table_args__ = (
        db.UniqueConstraint("degree", "semester", "name", name="uq_section_deg_sem_name"),
    )

class Subject(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    degree = db.Column(db.String(50), nullable=False)
    semester = db.Column(db.Integer, nullable=False)
    name = db.Column(db.String(200), nullable=False)
    __table_args__ = (
        db.UniqueConstraint("degree", "semester", "name", name="uq_subject_degree_sem_name"),
    )


class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    srn = db.Column(db.String(100), unique=True, nullable=True)
    name = db.Column(db.String(200), nullable=False)
    email = db.Column(db.String(200), unique=True, nullable=False)
    password_hash = db.Column(db.String(256), nullable=False)
    role = db.Column(db.String(20), default="student")
    degree = db.Column(db.String(50), nullable=True)
    semester = db.Column(db.Integer, nullable=True)
    section = db.Column(db.String(10), nullable=True)
    status = db.Column(db.String(20), default="PENDING")
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    emp_id = db.Column(db.String(100), unique=True, nullable=True)
    parent_email = db.Column(db.String(200), nullable=True)
    parent_password_hash = db.Column(db.String(256), nullable=True)


class Note(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(300))
    degree = db.Column(db.String(50))
    semester = db.Column(db.Integer)
    section = db.Column(db.String(50), nullable=True) # Added Section
    subject = db.Column(db.String(200))
    document_type = db.Column(db.String(100))
    file_path = db.Column(db.String(500))
    uploaded_by = db.Column(db.Integer, db.ForeignKey("user.id"))
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)


class Notice(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(300))
    message = db.Column(db.Text)
    degree = db.Column(db.String(50))
    semester = db.Column(db.Integer)
    section = db.Column(db.String(50))
    subject = db.Column(db.String(200))
    deadline = db.Column(db.DateTime, nullable=True)
    attachment = db.Column(db.String(500), nullable=True)
    professor_id = db.Column(db.Integer, db.ForeignKey("user.id"))
    professor_name = db.Column(db.String(200))
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class Book(db.Model):
    __tablename__ = "books"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    title = db.Column(db.String(300), nullable=False)
    author = db.Column(db.String(200))
    isbn = db.Column(db.String(100), nullable=True) # Added ISBN column for better searching
    degree = db.Column(db.String(50))
    semester = db.Column(db.Integer)
    file_path = db.Column(db.String(500), nullable=False)
    uploaded_by = db.Column(db.Integer, db.ForeignKey("user.id"))
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)


class HostelComplaint(db.Model):
    __tablename__ = "hostel_complaints"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    student_id = db.Column(db.Integer, db.ForeignKey("user.id"))
    hostel_id = db.Column(db.Integer, db.ForeignKey('hostel.id')) # NEW: Hostel info
    room_id = db.Column(db.Integer, db.ForeignKey('room.id')) 
    title = db.Column(db.String(300))
    description = db.Column(db.Text)
    status = db.Column(db.String(20), default="Open")
    attachment = db.Column(db.String(500), nullable=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    # NEW: Audit Trail/Status Log (JSON string)
    audit_trail = db.Column(db.Text, default='[]') 


class FeeNotification(db.Model):
    __tablename__ = "fee_notifications"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    title = db.Column(db.String(300))
    description = db.Column(db.Text)
    amount_cents = db.Column(db.BigInteger, nullable=False)
    category = db.Column(db.String(50), default="misc")
    issued_by = db.Column(db.Integer, db.ForeignKey("user.id"))
    due_date = db.Column(db.DateTime, nullable=True)
    target = db.Column(db.String(50), default="batch")
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class FeeTarget(db.Model):
    __tablename__ = "fee_targets"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    notification_id = db.Column(db.String, db.ForeignKey("fee_notifications.id", ondelete="CASCADE"))
    student_id = db.Column(db.Integer, db.ForeignKey("user.id"))
    status = db.Column(db.String(20), default="pending")
    paid_at = db.Column(db.DateTime, nullable=True)
    order_id = db.Column(db.String, nullable=True)


class Order(db.Model):
    __tablename__ = "orders"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    student_id = db.Column(db.Integer, nullable=False)
    amount_cents = db.Column(db.BigInteger, nullable=False)
    currency = db.Column(db.String(8), default="INR")
    description = db.Column(db.Text)
    status = db.Column(db.String(32), default="created")
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class Payment(db.Model):
    __tablename__ = "payments"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    order_id = db.Column(db.String, db.ForeignKey("orders.id", ondelete="cascade"))
    amount_cents = db.Column(db.BigInteger, nullable=False)
    payment_method = db.Column(db.String(64), nullable=True)
    transaction_id = db.Column(db.String(200), nullable=True)
    status = db.Column(db.String(32), nullable=False)
    extra_metadata = db.Column(db.JSON, default={})
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class Receipt(db.Model):
    __tablename__ = "receipts"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    payment_id = db.Column(db.String, db.ForeignKey("payments.id", ondelete="cascade"))
    storage_key = db.Column(db.String(500))
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class Mark(db.Model):
    __tablename__ = "marks"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    student_id = db.Column(db.Integer, db.ForeignKey("user.id"))
    subject = db.Column(db.String(200))
    exam_type = db.Column(db.String(50))
    marks_obtained = db.Column(db.Float)
    max_marks = db.Column(db.Float)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    uploaded_by = db.Column(db.Integer, db.ForeignKey("user.id"))


class Feedback(db.Model):
    __tablename__ = "feedback"
    id = db.Column(db.String, primary_key=True, default=lambda: str(uuid.uuid4()))
    student_id = db.Column(db.Integer, db.ForeignKey("user.id"))
    subject = db.Column(db.String(200))
    faculty_id = db.Column(db.Integer, db.ForeignKey("user.id"))
    text = db.Column(db.Text)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


class FacultyAllocation(db.Model):
    __tablename__ = "faculty_allocations"
    id = db.Column(db.Integer, primary_key=True)
    faculty_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    degree = db.Column(db.String(50), nullable=False)
    semester = db.Column(db.Integer, nullable=False)
    section = db.Column(db.String(50), nullable=False)
    subject = db.Column(db.String(200), nullable=False)
    # Unique constraint to prevent duplicate allocations
    __table_args__ = (
        db.UniqueConstraint("faculty_id", "degree", "semester", "section", "subject", name="_faculty_class_uc"),
    )


class Attendance(db.Model):
    __tablename__ = "attendance"
    id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    faculty_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    degree = db.Column(db.String(50), nullable=False)
    semester = db.Column(db.Integer, nullable=False)
    section = db.Column(db.String(50), nullable=False)
    subject = db.Column(db.String(200), nullable=False)
    date = db.Column(db.Date, nullable=False)
    status = db.Column(db.String(20), nullable=False) # "Present" or "Absent"
    timestamp = db.Column(db.DateTime, default=datetime.utcnow) # For 30 min edit window check
    __table_args__ = (
        db.UniqueConstraint("student_id", "subject", "date", name="_student_subject_date_uc"),
    )

# --- NEW PERSONAL ATTENDANCE MODELS ---
class StudentRoutine(db.Model):
    __tablename__ = 'student_routine'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    day_of_week = db.Column(db.String(10), nullable=False) # 'Monday', ...
    subjects = db.Column(db.Text, nullable=False) # JSON list
    __table_args__ = (
        db.UniqueConstraint('user_id', 'day_of_week', name='uq_std_routine_day'),
    )

class StudentAttendanceLog(db.Model):
    __tablename__ = 'student_attendance_log'
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    date = db.Column(db.Date, nullable=False)
    subject = db.Column(db.String(200), nullable=True) # NULL if "No Class"
    status = db.Column(db.String(20), nullable=False) # Present, Absent, No Class, Holiday
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)
    __table_args__ = (
        db.UniqueConstraint('user_id', 'date', 'subject', name='uq_std_log_date_sub'),
    )
# --- END PERSONAL ATTENDANCE MODELS ---


class Message(db.Model):
    __tablename__ = "messages"
    id = db.Column(db.Integer, primary_key=True)
    sender = db.Column(db.String(20), nullable=False) # 'parent' or 'faculty'
    student_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False)
    faculty_id = db.Column(db.Integer, db.ForeignKey("user.id"), nullable=False) 
    subject = db.Column(db.String(200))
    body = db.Column(db.Text)
    is_read = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)


# ==================== HRD (PLACEMENT CELL) MODELS ====================

class HRDUser(db.Model):
    __tablename__ = "hrd_users"
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(200), unique=True, nullable=False)
    password_hash = db.Column(db.String(256), nullable=False)
    name = db.Column(db.String(200), nullable=False)
    phone = db.Column(db.String(20), nullable=True)
    department = db.Column(db.String(100), default="Placement Cell")
    is_active = db.Column(db.Boolean, default=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class Company(db.Model):
    __tablename__ = "companies"
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    sector = db.Column(db.String(100), nullable=True)
    website = db.Column(db.String(300), nullable=True)
    hr_name = db.Column(db.String(200), nullable=True)
    hr_email = db.Column(db.String(200), nullable=True)
    hr_phone = db.Column(db.String(20), nullable=True)
    visit_history = db.Column(db.JSON, default=[])
    is_active = db.Column(db.Boolean, default=True)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class PlacementDrive(db.Model):
    __tablename__ = "placement_drives"
    id = db.Column(db.Integer, primary_key=True)
    company_id = db.Column(db.Integer, db.ForeignKey("companies.id", ondelete="CASCADE"), nullable=False)
    title = db.Column(db.String(300), nullable=False)
    description = db.Column(db.Text, nullable=True)
    role = db.Column(db.String(200), nullable=False)
    ctc_min = db.Column(db.Float, nullable=True) 
    ctc_max = db.Column(db.Float, nullable=True)
    location = db.Column(db.String(200), nullable=True)
    drive_type = db.Column(db.String(50), default="on_campus")
    eligibility_criteria = db.Column(db.JSON, default={})
    application_start = db.Column(db.DateTime, nullable=True)
    application_end = db.Column(db.DateTime, nullable=True)
    status = db.Column(db.String(50), default="open")
    created_by = db.Column(db.Integer, db.ForeignKey("hrd_users.id"), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class StudentPlacementProfile(db.Model):
    __tablename__ = "student_placement_profiles"
    id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey("user.id", ondelete="CASCADE"), unique=True, nullable=False)
    resume_url = db.Column(db.String(500), nullable=True)
    skills = db.Column(db.JSON, default=[])
    certifications = db.Column(db.JSON, default=[])
    projects = db.Column(db.JSON, default=[])
    github_url = db.Column(db.String(300), nullable=True)
    linkedin_url = db.Column(db.String(300), nullable=True)
    portfolio_url = db.Column(db.String(300), nullable=True)
    preferred_location = db.Column(db.String(200), nullable=True)
    preferred_role = db.Column(db.String(200), nullable=True)
    min_expected_ctc = db.Column(db.Float, nullable=True)
    profile_updated_at = db.Column(db.DateTime, default=datetime.utcnow)

class DriveApplication(db.Model):
    __tablename__ = "drive_applications"
    id = db.Column(db.Integer, primary_key=True)
    drive_id = db.Column(db.Integer, db.ForeignKey("placement_drives.id", ondelete="CASCADE"), nullable=False)
    student_id = db.Column(db.Integer, db.ForeignKey("user.id", ondelete="CASCADE"), nullable=False)
    applied_at = db.Column(db.DateTime, default=datetime.utcnow)
    status = db.Column(db.String(50), default="applied")
    withdrawal_reason = db.Column(db.Text, nullable=True)
    tags = db.Column(db.JSON, default=[])
    __table_args__ = (
        db.UniqueConstraint("drive_id", "student_id", name="uq_drive_application"),
    )

class InterviewRound(db.Model):
    __tablename__ = "interview_rounds"
    id = db.Column(db.Integer, primary_key=True)
    drive_id = db.Column(db.Integer, db.ForeignKey("placement_drives.id", ondelete="CASCADE"), nullable=False)
    round_name = db.Column(db.String(100), nullable=False)
    round_order = db.Column(db.Integer, default=1)
    scheduled_date = db.Column(db.DateTime, nullable=True)

class InterviewResult(db.Model):
    __tablename__ = "interview_results"
    id = db.Column(db.Integer, primary_key=True)
    round_id = db.Column(db.Integer, db.ForeignKey("interview_rounds.id", ondelete="CASCADE"), nullable=False)
    student_id = db.Column(db.Integer, db.ForeignKey("user.id", ondelete="CASCADE"), nullable=False)
    result = db.Column(db.String(50), nullable=False)
    feedback = db.Column(db.Text, nullable=True)
    recorded_by = db.Column(db.Integer, db.ForeignKey("hrd_users.id"), nullable=False)
    recorded_at = db.Column(db.DateTime, default=datetime.utcnow)

class PlacementOffer(db.Model):
    __tablename__ = "placement_offers"
    id = db.Column(db.Integer, primary_key=True)
    drive_id = db.Column(db.Integer, db.ForeignKey("placement_drives.id", ondelete="CASCADE"), nullable=False)
    student_id = db.Column(db.Integer, db.ForeignKey("user.id", ondelete="CASCADE"), nullable=False)
    role = db.Column(db.String(200), nullable=False)
    ctc = db.Column(db.Float, nullable=False)
    location = db.Column(db.String(200), nullable=True)
    offer_letter_url = db.Column(db.String(500), nullable=True)
    joining_date = db.Column(db.Date, nullable=True)
    bond_details = db.Column(db.Text, nullable=True)
    offer_date = db.Column(db.Date, default=datetime.utcnow)
    expiry_date = db.Column(db.Date, nullable=True)
    status = db.Column(db.String(50), default="pending")
    student_response_date = db.Column(db.DateTime, nullable=True)
    rejection_reason = db.Column(db.Text, nullable=True)

class HRDAnnouncement(db.Model):
    __tablename__ = "hrd_announcements"
    id = db.Column(db.Integer, primary_key=True)
    created_by = db.Column(db.Integer, db.ForeignKey("hrd_users.id"), nullable=False)
    title = db.Column(db.String(300), nullable=False)
    message = db.Column(db.Text, nullable=False)
    target_audience = db.Column(db.String(100), default="all")
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class AIAnalysisCache(db.Model):
    __tablename__ = "ai_analysis_cache"
    id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey("user.id", ondelete="CASCADE"), nullable=False)
    resume_url = db.Column(db.String(500), nullable=False)
    skills_extracted = db.Column(db.JSON, default=[])
    quality_score = db.Column(db.Integer, default=0)
    ats_score = db.Column(db.Integer, default=0)
    analyzed_at = db.Column(db.DateTime, default=datetime.utcnow)

class PlacementActivityLog(db.Model):
    __tablename__ = "placement_activity_logs"
    id = db.Column(db.Integer, primary_key=True)
    actor_id = db.Column(db.Integer, nullable=False)
    actor_role = db.Column(db.String(20), nullable=False)
    action = db.Column(db.String(100), nullable=False)
    entity_type = db.Column(db.String(50), nullable=False)
    entity_id = db.Column(db.Integer, nullable=False)
    details = db.Column(db.JSON, default={})
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)



def roles_allowed(roles):
    def decorator(fn):
        @wraps(fn)
        def wrapper(*args, **kwargs):
            # Verify JWT in request and return friendly JSON if missing/invalid
            try:
                verify_jwt_in_request()
            except Exception as e:
                # Return 401 and hint for client
                return jsonify({"success": False, "message": "Authentication required", "detail": str(e)}), 401

            try:
                identity = get_jwt_identity()
                user = User.query.get(identity)
            except Exception:
                return jsonify({"success": False, "message": "Invalid token or user not found"}), 401

            if not user or user.role not in roles:
                return jsonify({"success": False, "message": "Insufficient permissions"}), 403
            return fn(*args, **kwargs)
        return wrapper
    return decorator


def roles_allowed(allowed_roles):
    def decorator(fn):
        @wraps(fn)
        def wrapper(*args, **kwargs):
            try:
                verify_jwt_in_request()
                claims = get_jwt()
            except Exception as e:
                return jsonify({"success": False, "message": "Authentication required", "detail": str(e)}), 401

            if claims.get("role") not in allowed_roles:
                return jsonify({"success": False, "message": "Insufficient permissions"}), 403
        
            return fn(*args, **kwargs)
        return wrapper
    return decorator


def admin_only(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        try:
            # 1. VERIFY JWT and extract identity
            verify_jwt_in_request()
            claims = get_jwt()
        except Exception as e:
            return jsonify({"success": False, "message": "Authentication required", "detail": str(e)}), 401

        # 2. RBAC CHECK
        if claims.get("role") != "admin":
            return jsonify({"success": False, "message": "Admin access required"}), 403
        
        # 3. PASS CONTROL
        return fn(*args, **kwargs)
    return wrapper


# ==================== GROQ AI INTEGRATION FOR HRD ====================

def groq_ai_call(messages, temperature=0.7):
    """Make a request to Groq AI API using llama-3.3-70b-versatile model."""
    if not GROQ_API_KEY:
        return {"error": "GROQ_API_KEY not configured"}
    
    try:
        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama-3.3-70b-versatile",
                "messages": messages,
                "temperature": temperature,
                "max_tokens": 2048
            },
            timeout=30
        )
        response.raise_for_status()
        return response.json()
    except Exception as e:
        print(f"Groq AI Error: {e}")
        return {"error": str(e)}

def extract_skills_from_resume(resume_text):
    """Use Groq AI to extract skills from resume text."""
    messages = [
        {
            "role": "system",
            "content": "You are a resume parser. Extract skills, certifications, and projects from the resume. Return ONLY a valid JSON object with keys: skills (array of strings), certifications (array of strings), projects (array of strings). Do not include any markdown formatting or extra text."
        },
        {
            "role": "user",
            "content": f"Parse this resume:\n\n{resume_text[:3000]}"
        }
    ]
    
    result = groq_ai_call(messages, temperature=0.3)
    
    if "error" in result:
        return {"skills": [], "certifications": [], "projects": [], "error": result["error"]}
    
    try:
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "{}")
        content = content.strip()
        if content.startswith("```"):
            content = content.split("\n", 1)[1].rsplit("\n", 1)[0]
        parsed = json.loads(content)
        return {
            "skills": parsed.get("skills", []),
            "certifications": parsed.get("certifications", []),
            "projects": parsed.get("projects", [])
        }
    except Exception as e:
        print(f"Parse error: {e}")
        return {"skills": [], "certifications": [], "projects": [], "error": str(e)}

def analyze_resume_quality(resume_text):
    """Use Groq AI to score resume quality (0-100)."""
    messages = [
        {
            "role": "system",
            "content": "You are a resume quality analyzer. Rate the resume on a scale of 0-100 for both Quality and ATS compatibility. Return ONLY a valid JSON object with keys: quality_score (int 0-100), ats_score (int 0-100), feedback (string). Quality considers structure, grammar, achievements. ATS considers keyword density, formatting."
        },
        {
            "role": "user",
            "content": f"Analyze this resume:\n\n{resume_text[:3000]}"
        }
    ]
    
    result = groq_ai_call(messages, temperature=0.3)
    
    if "error" in result:
        return {"quality_score": 0, "ats_score": 0, "feedback": "Analysis failed", "error": result["error"]}
    
    try:
        content = result.get("choices", [{}])[0].get("message", {}).get("content", "{}")
        content = content.strip()
        if content.startswith("```"):
            content = content.split("\n", 1)[1].rsplit("\n", 1)[0]
        parsed = json.loads(content)
        return {
            "quality_score": int(parsed.get("quality_score", 0)),
            "ats_score": int(parsed.get("ats_score", 0)),
            "feedback": parsed.get("feedback", "")
        }
    except Exception as e:
        return {"quality_score": 0, "ats_score": 0, "feedback": "Parse error", "error": str(e)}

def calculate_role_fit(student_skills, job_requirements):
    """Calculate role-fit percentage and provide explainable insights."""
    student_set = set([s.lower().strip() for s in student_skills])
    job_set = set([r.lower().strip() for r in job_requirements])
    
    matched = list(student_set.intersection(job_set))
    missing = list(job_set - student_set)
    
    if len(job_set) == 0:
        match_pct = 0
    else:
        match_pct = int((len(matched) / len(job_set)) * 100)
    
    explanation = f"Match: {match_pct}%. Has: {', '.join(matched[:5]) if matched else 'None'}. Missing: {', '.join(missing[:5]) if missing else 'None'}"
    
    return {
        "match_percentage": match_pct,
        "matched": matched,
        "missing": missing,
        "explanation": explanation
    }

def analyze_batch_skills(student_skills_list):
    """Batch-level skill analysis."""
    skill_count = {}
    for skills in student_skills_list:
        for skill in skills:
            skill_lower = skill.lower()
            skill_count[skill_lower] = skill_count.get(skill_lower, 0) + 1
    
    sorted_skills = sorted(skill_count.items(), key=lambda x: x[1], reverse=True)
    top_skills = [{"skill": k, "count": v} for k, v in sorted_skills[:10]]
    
    return {
        "top_skills": top_skills,
        "total_unique_skills": len(skill_count),
        "distribution": skill_count
    }

# ==================== HRD HELPER DECORATORS ====================

def hrd_required(fn):
    """Decorator to ensure HRD role authentication"""
    @wraps(fn)
    def wrapper(*args, **kwargs):
        try:
            verify_jwt_in_request()
            claims = get_jwt()
        except Exception as e:
            return jsonify({"success": False, "message": "Authentication required"}), 401
        
        if claims.get("role") != "hrd":
            return jsonify({"success": False, "message": "HRD access required"}), 403
        
        return fn(*args, **kwargs)
    return wrapper

def log_hrd_activity(actor_id, actor_role, action, entity_type, entity_id, details={}):
    """Helper to log HRD activities"""
    log = PlacementActivityLog(
        actor_id=actor_id,
        actor_role=actor_role,
        action=action,
        entity_type=entity_type,
        entity_id=entity_id,
        details=details
    )
    db.session.add(log)
    db.session.commit()




# -------------------- AUTHENTICATION & SECURITY ENDPOINTS (OTP) --------------------


def mask_email(email):
    try:
        if "@" not in email: return email
        parts = email.split("@")
        user = parts[0]
        domain = parts[1]
        if len(user) <= 2:
            masked_user = user[0] + "***" 
        else:
            masked_user = user[:2] + "***" + user[-1]
        return f"{masked_user}@{domain}"
    except:
        return email

@app.route("/auth/send-otp", methods=["POST"])
def send_otp_endpoint():
    data = request.json or {}
    identifier = data.get("email") # Can be email or ID/SRN
    mode = data.get("mode", "signup")  # 'signup' or 'forgot_password'

    if not identifier:
        return jsonify({"success": False, "message": "Email or ID is required"}), 400

    email_to_use = identifier
    
    # Logic to resolve email if identifier is not an email (e.g. for Forgot Password)
    if mode == "forgot_password":
        if "@" not in identifier:
           user = User.query.filter((User.email == identifier) | (User.srn == identifier) | (User.emp_id == identifier)).first()
        else:
           user = User.query.filter_by(email=identifier).first()
           
        if not user:
            return jsonify({"success": False, "message": "User not found"}), 404
        email_to_use = user.email
    elif mode == "parent_forgot":
         # Parent Forgot: Input is SRN, Find Student, Use Parent Email
         user = User.query.filter_by(srn=identifier).first()
         if not user:
             return jsonify({"success": False, "message": "Student not found with this SRN"}), 404
         if not user.parent_email:
             return jsonify({"success": False, "message": "No parent email registered for this student."}), 400
         email_to_use = user.parent_email
    elif mode == "signup":
        # Check if user already exists
        if User.query.filter_by(email=email_to_use).first():
            return jsonify({"success": False, "message": "Email already registered. Please login."}), 400

    # Generate OTP
    otp_code = generate_otp()
    expires_at = datetime.utcnow() + timedelta(minutes=10)

    # Save to DB
    otp_record = OTP.query.filter_by(email=email_to_use).first()
    if otp_record:
        otp_record.otp = otp_code
        otp_record.expires_at = expires_at
        otp_record.verified = False
    else:
        otp_record = OTP(email=email_to_use, otp=otp_code, expires_at=expires_at, verified=False)
        db.session.add(otp_record)
    
    db.session.commit()
    print(f"DEBUG OTP for {email_to_use}: {otp_code}")

    # Send Email
    subject = "NoteOrbit - Verification Code"
    body = f"""
    <h2>NoteOrbit Verification</h2>
    <p>Your One-Time Password (OTP) is:</p>
    <h1 style="color: #3b82f6;">{otp_code}</h1>
    <p>This code expires in 10 minutes.</p>
    """
    
    success = send_email(email_to_use, subject, body)
    
    # Return masked email for UI confirmation
    masked = mask_email(email_to_use)
    
    if success:
        return jsonify({"success": True, "message": f"OTP sent to {masked}", "masked_email": masked, "email": email_to_use})
    else:
        # Fallback for dev/demo if SMTP fails but we want to confirm flow (based on debug log)
        return jsonify({"success": False, "message": "Failed to send email (Check server logs for OTP)", "masked_email": masked, "email": email_to_use}), 500


@app.route("/auth/lookup-user", methods=["POST"])
def lookup_user_endpoint():
    """Lookup user by ID/Email and return masked email for confirmation."""
    data = request.json or {}
    identifier = data.get("identifier", "").strip()
    
    if not identifier:
        return jsonify({"success": False, "message": "Identifier required"}), 400
        
    # Check email, SRN, or Emp ID
    user = User.query.filter((User.email == identifier) | (User.srn == identifier) | (User.emp_id == identifier)).first()
    
    if not user:
        return jsonify({"success": False, "message": "User not found"}), 404
        
    return jsonify({
        "success": True, 
        "masked_email": mask_email(user.email),
        "email_exists": True
    })

@app.route("/auth/lookup-parent", methods=["POST"])
def lookup_parent_endpoint():
    """Lookup parent by Ward's SRN and return masked email for confirmation."""
    data = request.json or {}
    srn = data.get("srn", "").strip()
    
    if not srn:
        return jsonify({"success": False, "message": "SRN required"}), 400
        
    # Find student by SRN
    student = User.query.filter_by(srn=srn).first()
    
    if not student:
        return jsonify({"success": False, "message": "Student not found with this SRN"}), 404
        
    if not student.parent_email:
        return jsonify({"success": False, "message": "No parent email registered for this student"}), 404
        
    return jsonify({
        "success": True, 
        "masked_email": mask_email(student.parent_email),
        "email_exists": True
    })

@app.route("/auth/verify-otp", methods=["POST"])
def verify_otp_endpoint():
    data = request.json or {}
    email = data.get("email", "").lower().strip() # Normalize email
    user_otp = data.get("otp")

    if not email or not user_otp:
        return jsonify({"success": False, "message": "Email and OTP required"}), 400

    otp_record = OTP.query.filter_by(email=email).first()
    
    if not otp_record:
        return jsonify({"success": False, "message": "No OTP request found"}), 400

    if otp_record.expires_at < datetime.utcnow():
        return jsonify({"success": False, "message": "OTP expired"}), 400

    if otp_record.otp != user_otp:
        return jsonify({"success": False, "message": "Invalid OTP"}), 400

    # Mark as verified
    otp_record.verified = True
    db.session.commit()

    return jsonify({"success": True, "message": "OTP verified successfully"})


@app.route("/auth/reset-password", methods=["POST"])
def reset_password_endpoint():
    data = request.json or {}
    email = data.get("email", "").lower().strip() # Normalize email
    otp = data.get("otp")
    new_password = data.get("new_password")
    role_type = data.get("role_type", "user") # 'user' or 'parent'

    if not email or not otp or not new_password:
        return jsonify({"success": False, "message": "Email, OTP, and new password required"}), 400

    # Verify OTP (Allow immediate verification if passed correctly)
    otp_record = OTP.query.filter_by(email=email).first()
    
    if not otp_record:
         return jsonify({"success": False, "message": "No OTP session found"}), 400
         
    # Check expiry
    if otp_record.expires_at < datetime.utcnow():
        return jsonify({"success": False, "message": "OTP expired"}), 400
        
    # Check content
    if otp_record.otp != otp:
        return jsonify({"success": False, "message": "Invalid OTP"}), 400
    
    # If we reached here, OTP is valid. No need to check 'verified' flag strictly
    # just treat it as verified since the correct OTP was provided in this request.

    # Update User Password
    if role_type == "parent":
        # Finds user by parent email (assuming unique parent email per student for this flow, or finding associated student via SRN if passed, but here we rely on email from OTP)
        user = User.query.filter_by(parent_email=email).first() 
    else:
        user = User.query.filter_by(email=email).first()

    if not user:
        return jsonify({"success": False, "message": "User not found"}), 404

    if role_type == "parent":
        user.parent_password_hash = hash_password(new_password)
    else:
        user.password_hash = hash_password(new_password)
        
    db.session.commit()

    # Clear OTP
    db.session.delete(otp_record)
    db.session.commit()

    return jsonify({"success": True, "message": "Password reset successfully. Please login."})


@app.route("/register", methods=["POST"])
def register():
    data = request.json or {}
    role = data.get("role")
    email = data.get("email")
    otp = data.get("otp")

    # Enforce OTP for students
    if role == "student":
        if not otp:
            return jsonify({"success": False, "message": "OTP is required for student registration"}), 400
        
        # Verify OTP record exists and is verified
        otp_record = OTP.query.filter_by(email=email).first()
        if not otp_record or not otp_record.verified or otp_record.otp != otp:
             return jsonify({"success": False, "message": "Email not verified. Please verify OTP first."}), 400

    if not email or not role:
        return jsonify({"success": False, "message": "Email and role required"}), 400

    if User.query.filter_by(email=email).first():
        return jsonify({"success": False, "message": "Email already registered"}), 400

    name = data.get("name")
    password = data.get("password")
    
    # Optional fields for student
    srn = data.get("srn")
    degree = data.get("degree")
    semester = data.get("semester")
    section = data.get("section")
    
    if srn and User.query.filter_by(srn=srn).first():
        return jsonify({"success": False, "message": "SRN already registered"}), 400

    new_user = User(
        email=email, 
        password_hash=hash_password(password), 
        role=role, 
        name=name,
        srn=srn, 
        degree=degree, 
        semester=int(semester) if semester else None, 
        section=section
    )
    
    if role == "student":
        new_user.status = "PENDING"
        msg = "Registration successful. Wait for Admin approval."
    else:
        new_user.status = "APPROVED"
        msg = "Faculty registration successful."

    db.session.add(new_user)
    
    # Cleanup OTP if used
    if otp:
        otp_record = OTP.query.filter_by(email=email).first()
        if otp_record:
            db.session.delete(otp_record)
            
    db.session.commit()
    return jsonify({"success": True, "message": msg})


@app.route("/login", methods=["POST"])
def login():
    data = request.json or {}
    identifier = data.get("email") # Can be email or SRN (for parents)
    password = data.get("password")
    role = data.get("role", "student") # student, admin, professor, parent

    if not identifier or not password:
        return jsonify({"success": False, "message": "Identifier & password required"}), 400

    if role == "parent":
        # Parent Login: uses SRN
        user = User.query.filter_by(srn=identifier).first()
        if not user:
             return jsonify({"success": False, "message": "Student not found"}), 404
        
        # Verify Parent Password
        if not user.parent_password_hash or user.parent_password_hash != hash_password(password):
             return jsonify({"success": False, "message": "Invalid parent credentials"}), 401

        # Success - Issue Token with parent role
        access_token = create_access_token(
            identity=str(user.id),
            additional_claims={"role": "parent"},
            expires_delta=timedelta(days=7)
        )
        return jsonify({
            "success": True,
            "token": access_token,
            "user": {
                "id": user.id, "name": f"{user.name}'s Parent", "email": user.parent_email, # Show Parent Email
                "role": "parent", "srn": user.srn,
                "degree": user.degree, "semester": user.semester, "section": user.section,
                "status": user.status
            }
        })

    # Normal Login (email based)
    user = User.query.filter_by(email=identifier).first()
    if not user:
        return jsonify({"success": False, "message": "User not found"}), 404

    if user.password_hash != hash_password(password):
        return jsonify({"success": False, "message": "Invalid credentials"}), 401

    if user.role == "student" and user.status != "APPROVED":
        return jsonify({"success": False, "message": f"Account status: {user.status}. Wait for admin approval."}), 403

    access_token = create_access_token(
        identity=str(user.id),
        additional_claims={"role": user.role},
        expires_delta=timedelta(days=7)
    )
    return jsonify({
        "success": True,
        "token": access_token,
        "user": {
            "id": user.id,
            "name": user.name,
            "email": user.email,
            "role": user.role,
            "degree": user.degree or "",
            "semester": user.semester or 1,
            "section": user.section or "",
            "status": user.status
        }
    })


@app.route("/me", methods=["GET"])
@jwt_required()
def me():
    identity = get_jwt_identity()
    user = User.query.get(identity)
    return jsonify({
        "id": user.id, "name": user.name, "email": user.email,
        "role": user.role, "degree": user.degree, "semester": user.semester,
        "section": user.section, "status": user.status
    })


@app.route("/admin/pending-students", methods=["GET"])
@admin_only
def pending_students():
    students = User.query.filter_by(role="student", status="PENDING").order_by(User.created_at.asc()).all()
    out = []
    for s in students:
        out.append({
            "id": s.id, "srn": s.srn, "name": s.name, "email": s.email,
            "degree": s.degree, "semester": s.semester, "section": s.section,
            "created_at": s.created_at.isoformat() if s.created_at else None
        })
    return jsonify({"success": True, "students": out}), 200


@app.route("/admin/approve-student", methods=["POST"])
@admin_only
def approve_student():
    data = request.json or {}
    student_id = data.get("student_id")
    action = data.get("action", "approve")
    try:
        sid = int(student_id)
    except Exception:
        return jsonify({"success": False, "message": "Invalid student_id"}), 400
    s = User.query.filter_by(id=sid, role="student").first()
    if not s:
        return jsonify({"success": False, "message": "Student not found"}), 404
    s.status = "APPROVED" if action == "approve" else "REJECTED"
    db.session.commit()
    try:
        if action == "approve":
            subject = "NoteOrbit - Account Approved"
            title = "Account Access Granted"
            main_body = f"Your account registration for NoteOrbit has been approved by the administrator. You can now log in to the portal and access all features."
            details = {
                "Name": s.name,
                "SRN": s.srn or "N/A",
                "Degree": s.degree or "N/A",
                "Semester": str(s.semester) if s.semester else "N/A",
                "Status": "Approved"
            }
            send_professional_email(s.email, subject, title, details, main_body)
        
        elif action == "reject":
            subject = "NoteOrbit - Account Rejected"
            title = "Account Registration Failed"
            main_body = f"Your account registration for NoteOrbit has been declined by the administrator. Please contact the administration if you believe this is an error."
            details = {
                "Name": s.name,
                "SRN": s.srn or "N/A",
                "Status": "Rejected"
            }
            send_professional_email(s.email, subject, title, details, main_body)
            
    except Exception as e:
        print(f"Error sending approval/rejection email: {e}")
        # Build success message even if email fails, but maybe note it? 
        # For now keeping it simple as per request.

    return jsonify({"success": True, "message": f"Student {action}d (Email notification sent)"})


@app.route("/admin/update-student", methods=["POST"])
@admin_only
def update_student():
    data = request.json or {}
    student_id = data.get("student_id")
    if not student_id:
        return jsonify({"success": False, "message": "Student ID required"}), 400
        
    s = User.query.get(student_id)
    if not s or s.role != "student":
        return jsonify({"success": False, "message": "Student not found"}), 404
        
    # Update allowed fields
    if "name" in data: s.name = data["name"]
    if "srn" in data: s.srn = data["srn"]
    if "degree" in data: s.degree = data["degree"]
    if "semester" in data: s.semester = int(data["semester"])
    if "section" in data: s.section = data["section"]
    if "parent_email" in data: s.parent_email = data["parent_email"]
    
    db.session.commit()
    return jsonify({"success": True, "message": "Student updated successfully"})


@app.route("/admin/generate-parent-password", methods=["POST"])
@admin_only
def generate_parent_password():
    data = request.json or {}
    student_id = data.get("student_id")
    if not student_id:
         return jsonify({"success": False, "message": "Student ID required"}), 400
         
    s = User.query.get(student_id)
    if not s or s.role != "student":
        return jsonify({"success": False, "message": "Student not found"}), 404
        
    if not s.parent_email:
        return jsonify({"success": False, "message": "Parent email not set for this student. Update student first."}), 400
        
    # Generate Random Password
    raw_password = f"P{random.randint(10000, 99999)}!"
    s.parent_password_hash = hash_password(raw_password)
    db.session.commit()
    
    # Send Email
    subject = f"NoteOrbit - Parent Portal Access for {s.name}"
    body = f"""
    <h2>Parent Portal Access Granted</h2>
    <p>Dear Parent,</p>
    <p>You have been granted access to the NoteOrbit Parent Portal to view the academic progress of your ward, <strong>{s.name}</strong> (SRN: {s.srn}).</p>
    <p><strong>Login Credentials:</strong></p>
    <ul>
        <li><strong>Ward SRN:</strong> {s.srn}</li>
        <li><strong>Password:</strong> {raw_password}</li>
    </ul>
    <p>Please login and change your password immediately.</p>
    <p><em>Access Link: <a href="http://localhost:5173/login">NoteOrbit Portal</a></em></p>
    """
    
    success = send_email(s.parent_email, subject, body)
    
    if success:
        return jsonify({"success": True, "message": f"Password sent to {s.parent_email}", "password": raw_password}) # Return pwd for dev ease if smtp fails
    else:
        return jsonify({"success": False, "message": "Failed to send email", "password": raw_password}), 500


@app.route("/admin/students", methods=["GET"])
@admin_only
def get_students_list_filtered():
    # Filters: degree, semester, section
    degree = request.args.get("degree")
    semester = request.args.get("semester")
    section = request.args.get("section")

    q = User.query.filter_by(role="student").order_by(User.name.asc())
    
    if degree:
        q = q.filter_by(degree=degree)
    
    if semester:
        try:
            q = q.filter_by(semester=int(semester))
        except ValueError:
            return jsonify({"success": False, "message": "Invalid semester parameter."}), 400
            
    if section:
        q = q.filter_by(section=section)
        
    students = q.all()
    out = []
    for s in students:
        # Fetch hostel/room info for display
        allocation = HostelAllocation.query.filter_by(student_id=s.id).first()
        hostel_info = None
        if allocation:
            hostel = Hostel.query.get(allocation.hostel_id)
            room = Room.query.get(allocation.room_id)
            if hostel and room:
                hostel_info = f"{hostel.name} (Room: {room.room_number})"
        
        out.append({
            "id": s.id, "srn": s.srn, "name": s.name, "email": s.email,
            "degree": s.degree, "semester": s.semester, "section": s.section,
            "status": s.status,
            "hostel_info": hostel_info, # Added hostel information
            "parent_email": s.parent_email,
            "parent_access": bool(s.parent_password_hash)
        })
        
    return jsonify({"success": True, "students": out}), 200


@app.route("/admin/add-faculty", methods=["POST"])
@admin_only
def add_faculty():
    data = request.json or {}
    required = ["name", "email", "password", "emp_id"]
    for r in required:
        if r not in data:
            return jsonify({"success": False, "message": f"Missing field {r}"}), 400

    if User.query.filter((User.email == data["email"]) | (User.emp_id == data["emp_id"])).first():
        return jsonify({"success": False, "message": "Email or EMP ID already registered."}), 400

    user = User(
        srn=None,
        name=data["name"],
        email=data["email"],
        password_hash=hash_password(data["password"]),
        role="professor",
        degree=None,
        semester=None,
        section=None,
        emp_id=data["emp_id"],
        status="APPROVED"
    )
    db.session.add(user)
    db.session.commit()
    return jsonify({"success": True, "message": f"Faculty account created for {user.name}."})


@app.route("/admin/faculty", methods=["GET"])
@admin_only
def list_faculty_allocations():
    faculty = User.query.filter_by(role="professor").all()
    out = []
    for f in faculty:
        allocs = FacultyAllocation.query.filter_by(faculty_id=f.id).all()
        alloc_data = [{
            "id": a.id, "degree": a.degree, "semester": a.semester, 
            "section": a.section, "subject": a.subject
        } for a in allocs]
        out.append({
            "id": f.id, "name": f.name, "email": f.email, 
            "emp_id": f.emp_id, "allocations": alloc_data
        })
    return jsonify({"success": True, "faculty": out})


@app.route("/admin/faculty/allocate", methods=["POST"])
@admin_only
def allocate_faculty():
    data = request.json or {}
    fid = data.get("faculty_id")
    degree = data.get("degree")
    semester = data.get("semester")
    section = data.get("section")
    subject = data.get("subject")

    if not (fid and degree and semester and section and subject):
        return jsonify({"success": False, "message": "All fields required"}), 400

    # Check for duplicate
    exists = FacultyAllocation.query.filter_by(
        faculty_id=fid, degree=degree, semester=int(semester), 
        section=section, subject=subject
    ).first()
    
    if exists:
        return jsonify({"success": False, "message": "Allocation already exists"}), 400

    # Create allocation
    new_alloc = FacultyAllocation(
        faculty_id=fid, degree=degree, semester=int(semester), 
        section=section, subject=subject
    )
    db.session.add(new_alloc)
    db.session.commit()
    return jsonify({"success": True, "message": "Class allocated to faculty."})


@app.route("/admin/faculty/deallocate", methods=["POST"])
@admin_only
def deallocate_faculty():
    data = request.json or {}
    id = data.get("allocation_id")
    alloc = FacultyAllocation.query.get(id)
    if not alloc:
        return jsonify({"success": False, "message": "Allocation not found"}), 404
        
    db.session.delete(alloc)
    db.session.commit()
    return jsonify({"success": True, "message": "Allocation removed."})


@app.route("/admin/degrees", methods=["GET", "POST", "DELETE"])
def manage_degrees():
    if request.method in ["POST", "DELETE"]:
        try:
            verify_jwt_in_request()
        except Exception:
            return jsonify({"success": False, "message": "Authentication required to modify degrees"}), 401
        user = User.query.get(get_jwt_identity())
        if user.role.lower() != "admin":
            return jsonify({"success": False, "message": "Admin access required to modify degrees"}), 403

    if request.method == "POST":
        name = (request.json or {}).get("name")
        if not name or Degree.query.filter_by(name=name).first():
            return jsonify({"success": False, "message": "Invalid or duplicate degree"}), 400
        db.session.add(Degree(name=name)); db.session.commit()

    if request.method == "DELETE":
        name = (request.json or {}).get("name")
        if not name:
            return jsonify({"success": False, "message": "Missing degree name for deletion"}), 400
        d = Degree.query.filter_by(name=name).first()
        if not d:
            return jsonify({"success": False, "message": "Degree not found"}), 404
        db.session.delete(d)
        db.session.commit()
        return jsonify({"success": True, "message": "Degree deleted successfully"})

    data = [d.name for d in Degree.query.order_by(Degree.name).all()]
    return jsonify({"success": True, "degrees": data})


@app.route("/admin/sections", methods=["GET", "POST", "DELETE"])
def manage_sections():
    # Helper to clean/int input
    def get_data():
        d = request.json or {}
        return d.get("degree"), d.get("semester"), d.get("name")

    if request.method in ["POST", "DELETE"]:
        try:
            verify_jwt_in_request()
        except Exception:
            return jsonify({"success": False, "message": "Authentication required"}), 401
        user = User.query.get(get_jwt_identity())
        if user.role != "admin":
            return jsonify({"success": False, "message": "Admin access required"}), 403

    if request.method == "POST":
        degree, semester, name = get_data()
        if not (degree and semester and name):
            return jsonify({"success": False, "message": "Missing fields"}), 400
        
        try:
            sem = int(semester)
        except:
             return jsonify({"success": False, "message": "Invalid semester"}), 400

        if Section.query.filter_by(degree=degree, semester=sem, name=name).first():
            return jsonify({"success": False, "message": "Duplicate section"}), 400
            
        db.session.add(Section(degree=degree, semester=sem, name=name))
        db.session.commit()
    
    if request.method == "DELETE":
        degree, semester, name = get_data()
        if not (degree and semester and name):
             return jsonify({"success": False, "message": "Missing fields for deletion"}), 400
             
        s = Section.query.filter_by(degree=degree, semester=int(semester), name=name).first()
        if not s:
            return jsonify({"success": False, "message": "Section not found"}), 404
        db.session.delete(s)
        db.session.commit()
        return jsonify({"success": True, "message": "Section deleted"})

    # GET: Filter by degree/sem query params
    degree = request.args.get("degree")
    semester = request.args.get("semester")
    
    q = Section.query
    if degree: q = q.filter_by(degree=degree)
    if semester:
        try:
            q = q.filter_by(semester=int(semester))
        except:
            pass
            
    # Return list of names? Or objects? 
    # Current frontend expects list of strings. But now sections are context-dependent.
    # If filtered, returning names is fine. If not filtered, returning unique names might be confusing 
    # but let's stick to returning names for the filtered context.
    
    data = [s.name for s in q.order_by(Section.name).all()]
    return jsonify({"success": True, "sections": data})

@app.route("/admin/subjects", methods=["GET", "POST", "DELETE"])
def manage_subjects():
    if request.method in ["POST", "DELETE"]:
        # === Manual Authorization & RBAC Check for WRITE requests ===
        try:
            verify_jwt_in_request()
            admin_claims = get_jwt()
        except Exception:
            return jsonify({"success": False, "message": "Authentication required to modify subjects"}), 401
            
        if admin_claims.get("role") not in ["admin", "professor"]:
            return jsonify({"success": False, "message": "Admin/Professor access required to modify subjects"}), 403
            
        payload = request.json or {}
        degree = payload.get("degree"); semester = payload.get("semester"); name = payload.get("name")
        
        if request.method == "POST":
            if not (degree and semester and name):
                return jsonify({"success": False, "message": "Missing fields"}), 400
            if not Degree.query.filter_by(name=degree).first():
                return jsonify({"success": False, "message": "Invalid degree"}), 400
            try:
                sem = int(semester)
            except:
                return jsonify({"success": False, "message": "Invalid semester"}), 400
            if Subject.query.filter_by(degree=degree, semester=sem, name=name).first():
                return jsonify({"success": False, "message": "Duplicate subject"}), 400
            db.session.add(Subject(degree=degree, semester=sem, name=name)); db.session.commit()

        if request.method == "DELETE":
            if not (degree and semester and name):
                return jsonify({"success": False, "message": "Missing fields for deletion"}), 400
            s = Subject.query.filter_by(degree=degree, semester=int(semester), name=name).first()
            if not s:
                return jsonify({"success": False, "message": "Subject not found"}), 404
            db.session.delete(s)
            db.session.commit()
            return jsonify({"success": True, "message": "Subject deleted successfully"})
    
    # --- GET logic (Read operations) ---
    degree = request.args.get("degree"); semester = request.args.get("semester"); q = Subject.query
    if degree: q = q.filter_by(degree=degree)
    if semester:
        try:
            q = q.filter_by(semester=int(semester))
        except:
            pass
    items = [{"degree": s.degree, "semester": s.semester, "name": s.name} for s in q.order_by(Subject.name).all()]
    return jsonify({"success": True, "subjects": items})


# -------------------- RESOURCES (Notes, Books, Notices) --------------------

@app.route("/upload-note", methods=["POST"])
@roles_allowed(["professor", "admin"])
def upload_note():
    title = request.form.get("title"); degree = request.form.get("degree")
    semester = request.form.get("semester"); section = request.form.get("section")
    subject = request.form.get("subject")
    document_type = request.form.get("document_type"); file = request.files.get("file")
    if not (title and degree and semester and section and subject and document_type and file):
        return jsonify({"success": False, "message": "Missing fields"}), 400
    if not allowed_file(file.filename):
        return jsonify({"success": False, "message": "File type not allowed"}), 400

    ext = file.filename.rsplit(".", 1)[1].lower()
    filename = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex}.{ext}"
    key = f"notes/{filename}"
    upload_to_minio(file, key, file.mimetype)

    identity = get_jwt_identity(); uploader = User.query.get(identity)
    note = Note(
        title=title, degree=degree, semester=int(semester), section=section, subject=subject,
        document_type=document_type, file_path=key, uploaded_by=uploader.id
    )
    db.session.add(note); db.session.commit()
    
    # Notify Students
    students = User.query.filter_by(role="student", degree=degree, semester=int(semester), section=section, status="APPROVED").all()
    for s in students:
        send_professional_email(
            s.email,
            f"New Note: {subject}",
            "New Study Material Uploaded",
            {"Subject": subject, "Title": title, "Type": document_type, "Faculty": uploader.name},
            f"New study material has been uploaded for <strong>{subject}</strong>."
        )
    
    return jsonify({"success": True, "message": "Note uploaded"})


@app.route("/notes", methods=["GET"])
@jwt_required()
def get_notes():
    user = User.query.get(get_jwt_identity())
    degree = request.args.get("degree") or user.degree; semester = request.args.get("semester") or user.semester
    section = request.args.get("section") or user.section
    subject = request.args.get("subject")
    document_type = request.args.get("document_type")
    if not degree or not semester:
        return jsonify({"success": False, "message": "degree and semester are required"}), 400
    try:
        sem = int(semester)
    except:
        return jsonify({"success": False, "message": "Invalid semester"}), 400
    q = Note.query.filter_by(degree=degree, semester=sem)
    if section:
        # Support "ALL" sections + User's specific section
        q = q.filter(or_(Note.section == section, Note.section == 'ALL'))

    if subject: q = q.filter_by(subject=subject)
    if document_type: q = q.filter_by(document_type=document_type)
    notes = q.order_by(Note.timestamp.desc()).all()
    out = []
    for n in notes:
        try:
            file_url = s3_client.generate_presigned_url(
                "get_object", Params={"Bucket": S3_BUCKET, "Key": n.file_path}, ExpiresIn=3600
            )
        except Exception:
            file_url = None
        out.append({
            "id": n.id, "title": n.title, "degree": n.degree, "semester": n.semester, "section": n.section,
            "subject": n.subject, "document_type": n.document_type, "file_url": file_url,
            "uploaded_by": n.uploaded_by, "timestamp": n.timestamp.isoformat() if n.timestamp else None
        })
    return jsonify({"success": True, "notes": out})

@app.route("/admin/notes/<int:note_id>", methods=["DELETE"])
@roles_allowed(["admin"])
def admin_delete_note(note_id):
    note = Note.query.get(note_id)
    if not note:
        return jsonify({"success": False, "message": "Note not found"}), 404

    # Best-effort delete from MinIO/S3 first (avoid orphaned objects)
    if note.file_path:
        try:
            s3_client.delete_object(Bucket=S3_BUCKET, Key=note.file_path)
        except Exception as e:
            # If object doesn't exist / transient error, still allow DB delete to proceed
            print(f"Failed to delete note object from storage: {e}")

    db.session.delete(note)
    db.session.commit()
    return jsonify({"success": True, "message": "Note deleted"})

# RENAMED and UPDATED for Admin Library Book Upload
@app.route("/api/admin/library/book", methods=["POST"])
@roles_allowed(["admin"])
def upload_book():
    title = request.form.get("title"); author = request.form.get("author")
    degree = request.form.get("degree"); semester = request.form.get("semester")
    isbn = request.form.get("isbn") # Optionally allow ISBN input
    file = request.files.get("file")
    
    if not (title and author and degree and semester and file):
        return jsonify({"success": False, "message": "Missing title, author, degree, semester, or file"}), 400
    if not allowed_file(file.filename):
        return jsonify({"success": False, "message": "File type not allowed"}), 400

    ext = file.filename.rsplit(".", 1)[1].lower()
    filename = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex}.{ext}"
    key = f"books/{filename}"
    upload_to_minio(file, key, file.mimetype)

    identity = get_jwt_identity()
    book = Book(
        title=title, author=author, isbn=isbn, degree=degree, semester=int(semester),
        file_path=key, uploaded_by=int(identity)
    )
    db.session.add(book); db.session.commit()
    return jsonify({"success": True, "message": f"Book '{title}' uploaded to Internal Library."})


# RENAMED and UPDATED for Unified Library Search
@app.route("/api/library/search", methods=["GET"])
@jwt_required()
def search_library():
    q = request.args.get("q", "")
    source = request.args.get("source", "internal").lower() # Default to internal
    
    if not q: return jsonify({"success": True, "books": []})
    
    out = []
    
    if source == "internal":
        # Internal Search Logic
        books = Book.query.filter(or_(
            Book.title.ilike(f"%{q}%"), 
            Book.author.ilike(f"%{q}%"),
            Book.isbn.ilike(f"%{q}%") # Search by ISBN too
        )).limit(20).all()

        for b in books:
            file_url = None
            try:
                # Generate a temporary download link
                file_url = s3_client.generate_presigned_url(
                    "get_object", Params={"Bucket": S3_BUCKET, "Key": b.file_path}, ExpiresIn=3600
                )
            except Exception:
                file_url = None
                
            # Internal book structure
            out.append({
                "id": b.id, "title": b.title, "author": b.author,
                "source": "Internal", "degree": b.degree, "semester": b.semester,
                "file_url": file_url,
                "isbn": b.isbn,
                "cover_url": None # Assuming no cover image is stored internally
            })
            
    elif source == "openlibrary":
        # External Search Logic
        out = search_open_library(q)
        
    else:
        return jsonify({"success": False, "message": "Invalid source parameter. Must be 'internal' or 'openlibrary'."}), 400
        
    return jsonify({"success": True, "books": out})


@app.route("/create-notice", methods=["POST"])
@roles_allowed(["professor", "admin"])
def create_notice():
    data = request.form or {}
    title = data.get("title"); message = data.get("message"); degree = data.get("degree")
    semester = data.get("semester"); section = data.get("section"); subject = data.get("subject")
    deadline_raw = data.get("deadline"); file = request.files.get("attachment")
    if not (title and message and degree and semester and section and subject):
        return jsonify({"success": False, "message": "Missing fields"}), 400
    if file and not allowed_file(file.filename):
        return jsonify({"success": False, "message": "Attachment type not allowed"}), 400

    attachment_key = None
    if file:
        ext = file.filename.rsplit(".", 1)[1].lower()
        fname = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex}.{ext}"
        attachment_key = f"notices/{fname}"
        upload_to_minio(file, attachment_key, file.mimetype)

    deadline = None
    if deadline_raw:
        try:
            deadline = datetime.strptime(deadline_raw, "%Y-%m-%d")
        except:
            return jsonify({"success": False, "message": "Invalid deadline format (YYYY-MM-DD)"}), 400

    identity = get_jwt_identity(); prof = User.query.get(identity)
    notice = Notice(
        title=title, message=message, degree=degree, semester=int(semester),
        section=section, subject=subject, deadline=deadline, attachment=attachment_key,
        professor_id=prof.id, professor_name=prof.name
    )
    db.session.add(notice); db.session.commit()

    # Notify Students
    try:
        sections_list = [s.strip().upper() for s in section.split(",")] if section else []
        q = User.query.filter_by(role="student", degree=degree, semester=int(semester), status="APPROVED")
        if sections_list:
             q = q.filter(User.section.in_(sections_list))
        
        students = q.all()
        for s in students:
             details = {"Subject": subject, "Posted By": prof.name, "Deadline": str(deadline) if deadline else "N/A"}
             send_professional_email(
                 s.email, 
                 f"New Notice: {title}", 
                 "Important Notice", 
                 details, 
                 f"{message[:200]}..." if len(message) > 200 else message
             )
    except Exception as e:
        print(f"Error sending notice emails: {e}")

    return jsonify({"success": True, "message": "Notice created"})


@app.route("/notices", methods=["GET"])
@jwt_required()
def get_notices():
    user = User.query.get(get_jwt_identity()); q = Notice.query
    subject = request.args.get("subject")

    if user.role == "student":
        q = q.filter_by(degree=user.degree, semester=user.semester)
        items = q.order_by(Notice.created_at.desc()).all(); filtered = []
        for n in items:
            sections = [s.strip().upper() for s in n.section.split(",")] if n.section else []
            if user.section and user.section.upper() in sections and (not subject or n.subject == subject):
                filtered.append(n)
        results = filtered
    else:
        degree = request.args.get("degree"); semester = request.args.get("semester"); section = request.args.get("section")
        if degree: q = q.filter_by(degree=degree)
        if semester:
            try:
                q = q.filter_by(semester=int(semester))
            except:
                return jsonify({"success": False, "message": "Invalid semester parameter."}), 400
        if section:
            q = q.filter(Notice.section.ilike(f"%{section}%"))
        if subject: q = q.filter_by(subject=subject)
        results = q.order_by(Notice.created_at.desc()).all()

    out = []
    for n in results:
        attachment_url = None
        if n.attachment:
            try:
                attachment_url = s3_client.generate_presigned_url(
                    "get_object", Params={"Bucket": S3_BUCKET, "Key": n.attachment}, ExpiresIn=3600
                )
            except Exception:
                attachment_url = None
        out.append({
            "id": n.id, "title": n.title, "message": n.message, "degree": n.degree,
            "semester": n.semester, "section": n.section, "subject": n.subject,
            "deadline": n.deadline.isoformat() if n.deadline else None, "attachment_url": attachment_url,
            "professor_id": n.professor_id, "professor_name": n.professor_name,
            "created_at": n.created_at.isoformat() if n.created_at else None
        })
    return jsonify({"success": True, "notices": out})


# -------------------- NEW HOSTEL MANAGEMENT ROUTES (ADMIN) --------------------

@app.route("/admin/hostel/hostels", methods=["GET", "POST"])
@admin_only
def manage_hostels():
    if request.method == "POST":
        data = request.json or {}
        name = data.get("name"); address = data.get("address")
        if not name:
            return jsonify({"success": False, "message": "Hostel name required"}), 400
        if Hostel.query.filter_by(name=name).first():
            return jsonify({"success": False, "message": "Hostel already exists"}), 400
        hostel = Hostel(name=name, address=address)
        db.session.add(hostel); db.session.commit()
        return jsonify({"success": True, "message": f"Hostel '{name}' added.", "id": hostel.id})
    
    # GET: Global Hostel Management View
    hostels = Hostel.query.all()
    out = []
    for h in hostels:
        total_rooms = Room.query.filter_by(hostel_id=h.id).count()
        occupied_rooms = Room.query.filter(Room.hostel_id == h.id, Room.current_occupancy > 0).count()
        total_capacity = db.session.query(func.sum(Room.capacity)).filter(Room.hostel_id == h.id).scalar() or 0
        current_occupancy = db.session.query(func.sum(Room.current_occupancy)).filter(Room.hostel_id == h.id).scalar() or 0
        
        out.append({
            "id": h.id, "name": h.name, "address": h.address,
            "total_rooms": total_rooms, "occupied_rooms": occupied_rooms,
            "total_capacity": total_capacity, "current_occupancy": current_occupancy,
            "vacant_beds": total_capacity - current_occupancy
        })
    return jsonify({"success": True, "hostels": out})


@app.route("/admin/hostel/rooms", methods=["GET", "POST"])
@admin_only
def manage_rooms():
    if request.method == "POST":
        data = request.json or {}
        hostel_id = data.get("hostel_id"); room_number = data.get("room_number")
        capacity = int(data.get("capacity", 1))
        if not (hostel_id and room_number) or capacity < 1:
            return jsonify({"success": False, "message": "Missing fields or invalid capacity"}), 400
        hostel = Hostel.query.get(hostel_id)
        if not hostel:
            return jsonify({"success": False, "message": "Hostel not found"}), 404
        if Room.query.filter_by(hostel_id=hostel_id, room_number=room_number).first():
            return jsonify({"success": False, "message": "Room already exists in this hostel"}), 400
        
        room = Room(hostel_id=hostel_id, room_number=room_number, capacity=capacity, current_occupancy=0)
        db.session.add(room); db.session.commit()
        # Update hostel total_rooms count (optional optimization)
        Hostel.query.filter_by(id=hostel_id).update({"total_rooms": Hostel.total_rooms + 1})
        db.session.commit()
        
        return jsonify({"success": True, "message": f"Room {room_number} added to {hostel.name}.", "id": room.id})

    # GET: Rooms List/Search
    hostel_id = request.args.get("hostel_id")
    q = Room.query
    if hostel_id:
        q = q.filter_by(hostel_id=hostel_id)
    rooms = q.order_by(Room.hostel_id, Room.room_number).all()
    
    out = []
    for r in rooms:
        hostel_name = Hostel.query.get(r.hostel_id).name if r.hostel_id else "N/A"
        out.append({
            "id": r.id, "hostel_id": r.hostel_id, "hostel_name": hostel_name,
            "room_number": r.room_number, "capacity": r.capacity, "occupancy": r.current_occupancy,
            "is_vacant": r.current_occupancy < r.capacity
        })
    return jsonify({"success": True, "rooms": out})


# MODIFIED: Accepts SRN instead of student_id in the payload
@app.route("/admin/hostel/assign-room", methods=["POST"])
@admin_only
def assign_room_to_student():
    data = request.json or {}
    # Retrieve SRN from the payload
    srn = data.get("srn"); room_id = data.get("room_id")
    
    if not (srn and room_id):
        return jsonify({"success": False, "message": "Missing SRN or room_id"}), 400

    # NEW STEP: Find student by SRN
    student = User.query.filter_by(srn=srn).first()
    
    if not student or student.role != 'student':
        return jsonify({"success": False, "message": "Student not found or user is not a student"}), 404
    
    # Use the retrieved student's internal ID for core logic
    student_id = student.id
    
    room = Room.query.get(room_id)
    
    if not room:
        return jsonify({"success": False, "message": "Room not found"}), 404
    if room.current_occupancy >= room.capacity:
        return jsonify({"success": False, "message": "Room is fully occupied"}), 400

    # 1. Remove previous allocation if exists (Core Logic Start)
    existing_allocation = HostelAllocation.query.filter_by(student_id=student_id).first()
    if existing_allocation:
        # Decrement occupancy in the old room
        old_room = Room.query.get(existing_allocation.room_id)
        if old_room:
            old_room.current_occupancy = max(0, old_room.current_occupancy - 1)
        db.session.delete(existing_allocation)
    
    # 2. Create new allocation
    new_allocation = HostelAllocation(
        student_id=student_id,
        hostel_id=room.hostel_id,
        room_id=room_id
    )
    db.session.add(new_allocation)
    
    # 3. Increment occupancy in the new room
    room.current_occupancy += 1
    
    db.session.commit()
    # Core Logic End
    return jsonify({"success": True, "message": f"Room {room.room_number} assigned to {student.name} (SRN: {srn})."})

# -------------------- HOSTEL COMPLAINT ROUTES (MODIFIED) --------------------

@app.route("/hostel/complaints", methods=["POST"])
@roles_allowed(["student"])
def submit_hostel_complaint():
    student_id = int(get_jwt_identity())
    
    # NEW: Check for hostel allocation before allowing complaint
    allocation = HostelAllocation.query.filter_by(student_id=student_id).first()
    if not allocation:
        # User is not allocated a hostel, block complaint submission
        return jsonify({"success": False, "message": "Complaint submission not allowed. No hostel is currently allotted to your account."}), 403

    title = request.form.get("title"); description = request.form.get("description"); file = request.files.get("attachment")
    if not title or not description:
        return jsonify({"success": False, "message": "Title and description required"}), 400

    attachment_key = None
    if file:
        if not allowed_file(file.filename):
            return jsonify({"success": False, "message": "Attachment file type not allowed."}), 400
        ext = file.filename.rsplit(".", 1)[1].lower()
        fname = f"{datetime.utcnow().strftime('%Y%m%d%H%M%S')}_{uuid.uuid4().hex}.{ext}"
        attachment_key = f"hostel-complaints/{fname}"
        upload_to_minio(file, attachment_key, file.mimetype)

    # Initial audit trail log
    initial_trail = json.dumps([{"status": "Open", "timestamp": datetime.utcnow().isoformat(), "note": "Complaint submitted.", "by": "Student"}])
    
    complaint = HostelComplaint(
        student_id=student_id, 
        hostel_id=allocation.hostel_id, # New: Store hostel info
        room_id=allocation.room_id, 
        title=title, 
        description=description, 
        attachment=attachment_key,
        status="Open",
        audit_trail=initial_trail
    )
    db.session.add(complaint); db.session.commit()
    return jsonify({"success": True, "message": "Complaint submitted successfully. Status: Open"})


# NEW ROUTE: Student Portal - View My Complaints and History
@app.route("/student/hostel/complaints", methods=["GET"])
@roles_allowed(["student", "parent"])
def student_view_hostel_complaints():
    student_id = int(get_jwt_identity())
    
    # Query for complaints belonging to the current student
    complaints = db.session.query(HostelComplaint, Hostel, Room).\
        filter(HostelComplaint.student_id == student_id).\
        outerjoin(Hostel, HostelComplaint.hostel_id == Hostel.id).\
        outerjoin(Room, HostelComplaint.room_id == Room.id).\
        order_by(HostelComplaint.created_at.desc()).all()

    out = []
    for c, hostel, room in complaints:
        file_url = None
        if c.attachment:
            try:
                # Generate a temporary download link for the attachment
                file_url = s3_client.generate_presigned_url(
                    "get_object", Params={"Bucket": S3_BUCKET, "Key": c.attachment}, ExpiresIn=3600
                )
            except Exception:
                file_url = None
                
        # Parse audit trail for live tracking/history
        try:
            audit_log = json.loads(c.audit_trail)
        except:
            audit_log = [{"status": c.status, "timestamp": c.created_at.isoformat() if c.created_at else "N/A", "note": "Initial status"}]
            
        out.append({
            "id": c.id, 
            "title": c.title, 
            "description": c.description,
            "status": c.status, 
            "hostel_name": hostel.name if hostel else "N/A", 
            "room_number": room.room_number if room else "N/A", 
            "file_url": file_url,
            "audit_trail": audit_log, # This provides the history/live tracking
            "created_at": c.created_at.isoformat() if c.created_at else None
        })
    return jsonify({"success": True, "complaints": out})
# END NEW ROUTE


@app.route("/admin/hostel/complaints", methods=["GET"])
@admin_only
def view_hostel_complaints():
    complaints = db.session.query(HostelComplaint, User, Hostel, Room).\
        join(User, HostelComplaint.student_id == User.id).\
        outerjoin(Hostel, HostelComplaint.hostel_id == Hostel.id).\
        outerjoin(Room, HostelComplaint.room_id == Room.id).\
        order_by(HostelComplaint.created_at.desc()).all()
        
    out = []
    for c, student, hostel, room in complaints:
        file_url = None
        if c.attachment:
            try:
                file_url = s3_client.generate_presigned_url(
                    "get_object", Params={"Bucket": S3_BUCKET, "Key": c.attachment}, ExpiresIn=3600
                )
            except Exception:
                file_url = None
                
        # Parse audit trail for frontend display
        try:
            audit_log = json.loads(c.audit_trail)
        except:
            audit_log = [{"status": c.status, "timestamp": c.created_at.isoformat() if c.created_at else "N/A", "note": "Initial status"}]
            
        out.append({
            "id": c.id, "title": c.title, "description": c.description,
            "status": c.status, 
            "student_id": student.id,
            "student_name": student.name if student else "Unknown",
            "hostel_name": hostel.name if hostel else "N/A", 
            "room_number": room.room_number if room else "N/A", 
            "file_url": file_url,
            "audit_trail": audit_log, # New: Audit Trail
            "created_at": c.created_at.isoformat() if c.created_at else None
        })
    return jsonify({"success": True, "complaints": out})


@app.route("/admin/hostel/complaints/<complaint_id>/status", methods=["PATCH"])
@admin_only
def update_hostel_complaint_status(complaint_id):
    data = request.json or {}
    new_status = data.get("status")
    note = data.get("note", "")
    
    VALID_STATUSES = ["Open", "Under Review", "Under Progress", "Resolved", "Closed"]
    
    if new_status not in VALID_STATUSES:
        return jsonify({"success": False, "message": "Invalid status value"}), 400
        
    complaint = HostelComplaint.query.get(complaint_id)
    if not complaint:
        return jsonify({"success": False, "message": "Complaint not found"}), 404
        
    # Update status
    complaint.status = new_status
    
    # Update audit trail
    trail = json.loads(complaint.audit_trail)
    admin_user = User.query.get(get_jwt_identity())
    trail.append({
        "status": new_status,
        "timestamp": datetime.utcnow().isoformat(),
        "by": admin_user.name if admin_user else "Admin",
        "note": note
    })
    complaint.audit_trail = json.dumps(trail)
    db.session.commit()
    
    # Notify Student
    s = User.query.get(complaint.student_id)
    if s:
        details = {"Complaint ID": f"#{complaint.id[:8]}", "Title": complaint.title, "New Status": new_status, "Admin Note": note}
        send_professional_email(s.email, f"Complaint Update: {new_status}", "Hostel Complaint Updated", details, f"The status of your hostel complaint '<strong>{complaint.title}</strong>' has been updated.")
    
    return jsonify({"success": True, "message": f"Complaint status updated to {new_status}.", "new_status": new_status})


@app.route("/admin/hostel/resolve", methods=["POST"])
@admin_only
def resolve_hostel_complaint():
    # Helper endpoint for frontend legacy button to use the new PATCH logic
    data = request.json or {}
    complaint_id = data.get("complaint_id")
    if not complaint_id:
        return jsonify({"success": False, "message": "Missing complaint_id"}), 400
    # Use the new PATCH logic internally to ensure audit trail is updated
    request.json = {"status": "Resolved", "note": "Complaint resolved by Admin using shortcut."}
    return update_hostel_complaint_status(complaint_id)


# -------------------- FEES, MARKS, FEEDBACK & AI CHAT ROUTES --------------------

@app.route("/admin/fees/create", methods=["POST"])
@admin_only
def admin_create_fee_notification():
    payload = request.json or {}
    title = payload.get("title"); amount_cents = int(payload.get("amount_cents", 0))
    if not title or amount_cents <= 0:
        return jsonify({"success": False, "message": "title and positive amount_cents required"}), 400

    notif = FeeNotification(
        title=title, description=payload.get("description"), amount_cents=amount_cents,
        category=payload.get("category", "misc"), issued_by=get_jwt_identity(),
        due_date=(datetime.strptime(payload["due_date"], "%Y-%m-%d") if payload.get("due_date") else None),
        target=payload.get("target", "batch")
    )
    db.session.add(notif); db.session.flush()

    target = notif.target; created_targets = 0
    if target in ("batch", "sem"):
        degree = payload.get("degree"); semester = payload.get("semester"); sections = payload.get("sections", "")
        q = User.query.filter_by(role="student")
        if degree: q = q.filter_by(degree=degree)
        if semester: q = q.filter_by(semester=int(semester))
        if sections:
            secs = [s.strip().upper() for s in sections.split(",") if s.strip()]
            students = q.filter(User.section.in_(secs)).all()
        else:
            students = q.all()
        for s in students:
            ft = FeeTarget(notification_id=notif.id, student_id=s.id)
            db.session.add(ft); created_targets += 1
            # Notify
            details = {"Title": title, "Amount": f"INR {amount_cents/100}", "Due Date": str(payload.get("due_date"))}
            send_professional_email(s.email, f"Fee Demand: {title}", "New Fee Notification", details, "A new fee payment is due.")
            if s.parent_email:
                send_professional_email(s.parent_email, f"Fee Demand: {s.name}", f"Fee Notification for {s.name}", details, f"A new fee payment is requested for your ward.")
    elif target == "custom":
        srns = payload.get("srns", []) or []
        for srn in srns:
            user = User.query.filter_by(srn=srn).first()
            if user:
                ft = FeeTarget(notification_id=notif.id, student_id=user.id)
                db.session.add(ft); created_targets += 1
                # Notify
                details = {"Title": title, "Amount": f"INR {amount_cents/100}", "Due Date": str(payload.get("due_date"))}
                send_professional_email(user.email, f"Fee Demand: {title}", "New Fee Notification", details, "A new fee payment is due.")
                if user.parent_email:
                    send_professional_email(user.parent_email, f"Fee Demand: {user.name}", f"Fee Notification for {user.name}", details, f"A new fee payment is requested for your ward.")
    elif target == "single":
        srn = payload.get("single_srn")
        user = User.query.filter_by(srn=srn).first()
        if user:
            ft = FeeTarget(notification_id=notif.id, student_id=user.id)
            db.session.add(ft); created_targets += 1
            # Notify
            details = {"Title": title, "Amount": f"INR {amount_cents/100}", "Due Date": str(payload.get("due_date"))}
            send_professional_email(user.email, f"Fee Demand: {title}", "New Fee Notification", details, "A new fee payment is due.")
            if user.parent_email:
                send_professional_email(user.parent_email, f"Fee Demand: {user.name}", f"Fee Notification for {user.name}", details, f"A new fee payment is requested for your ward.")
    db.session.commit()
    return jsonify({"success": True, "notification_id": notif.id, "targets_created": created_targets})


@app.route("/fees/list", methods=["GET"])
@jwt_required()
def student_fees_list():
    uid = int(get_jwt_identity())
    results = db.session.query(FeeTarget, FeeNotification).join(FeeNotification, FeeTarget.notification_id == FeeNotification.id).filter(FeeTarget.student_id == uid).order_by(FeeNotification.created_at.desc()).all()
    items = []
    for ft, notif in results:
        items.append({
            "target_id": ft.id, "notification_id": notif.id, "title": notif.title,
            "description": notif.description, "amount_cents": notif.amount_cents,
            "amount": cents_to_rupees_str(notif.amount_cents), "category": notif.category,
            "due_date": notif.due_date.isoformat() if notif.due_date else None,
            "status": ft.status, "paid_at": ft.paid_at.isoformat() if ft.paid_at else None,
            "payment_id": ft.order_id or None
        })
    return jsonify({"success": True, "fees": items})


@app.route("/fees/pay", methods=["POST"])
@jwt_required()
def create_demo_order():
    payload = request.json or {}; target_id = payload.get("target_id")
    ft = FeeTarget.query.filter_by(id=target_id).first()
    if not ft:
        return jsonify({"success": False, "message": "Target not found"}), 404
    notif = FeeNotification.query.get(ft.notification_id)
    if not notif:
        return jsonify({"success": False, "message": "Notification not found"}), 404
    if ft.status == "paid":
        return jsonify({"success": False, "message": "Fee already paid"}), 400
    order = Order(
        student_id=int(get_jwt_identity()), amount_cents=notif.amount_cents, description=notif.title
    )
    db.session.add(order); db.session.commit()
    ft.order_id = order.id; db.session.commit()
    return jsonify({"success": True, "order_id": order.id, "amount_cents": order.amount_cents})


@app.route("/demo/checkout/<order_id>", methods=["GET", "POST"])
def demo_checkout(order_id):
    order = Order.query.get(order_id)
    if not order:
        return "Order not found", 404
    if request.method == "GET":
        return f"""
        <html><body style="font-family: system-ui, Arial; padding: 24px;">
        <h2>Demo Checkout — NoteOrbit</h2><p>Order: {order.id} — Amount: ₹{cents_to_rupees_str(order.amount_cents)}</p>
        <form method="POST"><label>Card number (fake): <input name="card" /></label><br/><br/>
        <label>Name on card: <input name="name" /></label><br/><br/><button type="submit">Pay (Demo)</button></form>
        </body></html>
        """
    form = request.form
    payment = Payment(
        order_id=order.id, amount_cents=order.amount_cents, payment_method="card-demo",
        transaction_id=f"DEMO-{uuid.uuid4()}", status="success", extra_metadata={"card": form.get("card")}
    )
    db.session.add(payment); order.status = "paid"; db.session.commit()

    ft = FeeTarget.query.filter_by(order_id=order.id).first()
    if ft:
        ft.status = "paid";
        ft.paid_at = datetime.utcnow();
        ft.order_id = payment.id
        db.session.commit()

    student = User.query.get(order.student_id)
    student_name = student.name if student else "Student"; student_srn = student.srn if student else "SRN-NA"
    receipt_id = str(uuid.uuid4())
    html = f"""
    <html><body><h2>NoteOrbit — Payment Receipt (Demo)</h2><p><strong>Receipt ID:</strong> {receipt_id}</p>
    <p><strong>Student:</strong> {student_name} ({student_srn})</p><p><strong>Order ID:</strong> {order.id}</p>
    <p><strong>Payment ID:</strong> {payment.id}</p>
    <p><strong>Amount:</strong> ₹{cents_to_rupees_str(order.amount_cents)}</p><p><strong>Txn ID:</strong> {payment.transaction_id}</p>
    <p><small>Demo receipt generated by NoteOrbit. No real money exchanged.</small></p></body></html>
    """
    filename_base = f"receipt_{payment.id}"
    local_pdf = ensure_receipt_pdf(html, filename_base)
    storage_key, public_url = upload_receipt(local_pdf, dest_key=RECEIPT_PREFIX + os.path.basename(local_pdf))
    receipt = Receipt(payment_id=payment.id, storage_key=storage_key)
    db.session.add(receipt); db.session.commit()
    return f"""
    <html><body><h2>Payment Success (Demo)</h2><p>Transaction id: {payment.transaction_id}</p>
    <p><a href="{public_url}" target="_blank">Download Receipt (valid 1 hour)</a></p></body></html>
    """


@app.route("/fees/receipt/<payment_id>", methods=["GET"])
@jwt_required()
def get_receipt_for_payment(payment_id):
    pay = Payment.query.get(payment_id)
    if not pay:
        return jsonify({"success": False, "message": "Payment not found"}), 404
    rec = Receipt.query.filter_by(payment_id=payment_id).first()
    if not rec:
        return jsonify({"success": False, "message": "Receipt not found"}), 404
    url = s3_client.generate_presigned_url(
        "get_object", Params={"Bucket": S3_BUCKET, "Key": rec.storage_key}, ExpiresIn=3600
    )
    return jsonify({"success": True, "receipt_url": url})


@app.route("/faculty/marks/upload", methods=["POST"])
@roles_allowed(["professor", "admin"])
def faculty_upload_marks():
    payload = request.json or {}
    subject = payload.get("subject"); exam_type = payload.get("exam_type")
    srn = payload.get("srn")
    marks_obtained = payload.get("marks_obtained")
    max_marks = payload.get("max_marks")
    
    if not (subject and exam_type and srn and marks_obtained is not None and max_marks is not None):
        return jsonify({"success": False, "message": "subject, exam_type, srn, marks_obtained, and max_marks required"}), 400
        
    user = User.query.filter_by(srn=srn).first()
    if not user:
        return jsonify({"success": False, "message": f"Student with SRN {srn} not found."}), 404
        
    try:
        marks = float(marks_obtained)
        max_m = float(max_marks)
        if marks < 0 or max_m <= 0 or marks > max_m:
             return jsonify({"success": False, "message": "Invalid mark values (marks must be $\\ge 0$ and $\\le max\\_marks$, max\\_marks $> 0$)."}), 400
    except ValueError:
        return jsonify({"success": False, "message": "Marks must be numeric."}), 400

    uploader_id = int(get_jwt_identity())
    
    m = Mark(
        student_id=user.id, subject=subject, exam_type=exam_type,
        marks_obtained=marks, max_marks=max_m,
        uploaded_by=uploader_id
    )
    db.session.add(m); 
    db.session.commit()
    
    # Notify Student
    details = {"Subject": subject, "Exam Type": exam_type, "Score": f"{marks}/{max_m}", "Percentage": f"{(marks/max_m)*100:.1f}%"}
    send_professional_email(
        user.email, 
        f"Marks Released: {subject}", 
        "New Assessment Score", 
        details, 
        "Your marks for the recent assessment have been published."
    )
    
    # Notify Parent
    if user.parent_email:
        send_professional_email(
            user.parent_email, 
            f"Academic Alert: {user.name} - {subject}", 
            f"Marks Published for {user.name}", 
            details, 
            f"New marks have been uploaded for your ward, <strong>{user.name}</strong>."
        )
    
    return jsonify({"success": True, "message": f"Marks uploaded for {user.name}."})


# Removed the legacy /faculty/marks/upload-csv route as requested

@app.route("/student/marks", methods=["GET"])
@roles_allowed(["student", "parent"])
def student_get_marks():
    uid = int(get_jwt_identity()); rows = Mark.query.filter_by(student_id=uid).order_by(Mark.created_at.desc()).all()
    out = {}
    for r in rows:
        subj = r.subject
        out.setdefault(subj, []).append({
            "exam_type": r.exam_type,
            "marks_obtained": r.marks_obtained,
            "max_marks": r.max_marks,
            "uploaded_at": r.created_at.isoformat() if r.created_at else None
        })
    return jsonify({"success": True, "marks": out})


@app.route("/faculty/feedback", methods=["POST"])
@roles_allowed(["professor", "admin"])
def faculty_add_feedback():
    payload = request.json or {}; srn = payload.get("srn"); subject = payload.get("subject"); text = payload.get("text")
    if not srn or not subject or not text:
        return jsonify({"success": False, "message": "srn, subject and text required"}), 400
    student = User.query.filter_by(srn=srn).first()
    if not student:
        return jsonify({"success": False, "message": "student not found"}), 404
    fb = Feedback(
        student_id=student.id, subject=subject, faculty_id=int(get_jwt_identity()), text=text
    )
    db.session.add(fb); db.session.commit(); 
    
    # Notify Student & Parent
    details = {"Subject": subject, "Faculty": User.query.get(int(get_jwt_identity())).name, "Feedback": text}
    send_professional_email(student.email, f"New Feedback: {subject}", "Faculty Feedback Received", details, "You have received new feedback from your professor.")
    if student.parent_email:
        send_professional_email(student.parent_email, f"Feedback: {student.name} - {subject}", f"Faculty Feedback for {student.name}", details, f"Faculty has provided feedback for your ward, <strong>{student.name}</strong>.")

    return jsonify({"success": True, "message": "Feedback saved"})


@app.route("/student/feedback", methods=["GET"])
@roles_allowed(["student", "parent"])
def student_get_feedback():
    uid = int(get_jwt_identity()); rows = Feedback.query.filter_by(student_id=uid).order_by(Feedback.created_at.desc()).all()
    out = [{"subject": r.subject, "text": r.text, "faculty_id": r.faculty_id, "created_at": r.created_at.isoformat() if r.created_at else None} for r in rows]
    return jsonify({"success": True, "feedback": out})


@app.route("/faculty/allocations", methods=["GET"])
@roles_allowed(["professor", "admin"])
def get_faculty_allocations():
    fid = int(get_jwt_identity())
    allocs = FacultyAllocation.query.filter_by(faculty_id=fid).all()
    out = [{
        "id": a.id, "degree": a.degree, "semester": a.semester, 
        "section": a.section, "subject": a.subject
    } for a in allocs]
    return jsonify({"success": True, "allocations": out})


@app.route("/faculty/students", methods=["GET"])
@roles_allowed(["professor", "admin"])
def get_students_for_marking():
    degree = request.args.get("degree")
    semester = request.args.get("semester")
    section = request.args.get("section")
    subject = request.args.get("subject")
    date_str = request.args.get("date")

    if not (degree and semester and section):
        return jsonify({"success": False, "message": "Missing params"}), 400
    
    students = User.query.filter_by(role="student", degree=degree, semester=int(semester), section=section, status="APPROVED").order_by(User.srn.asc()).all()
    
    att_map = {}
    if subject and date_str:
        try:
            d = datetime.strptime(date_str, "%Y-%m-%d").date()
            rows = Attendance.query.filter_by(subject=subject, date=d).all()
            for r in rows:
                att_map[r.student_id] = {"status": r.status, "timestamp": r.timestamp}
        except ValueError:
            pass            

    out = []
    for s in students:
        att = att_map.get(s.id)
        out.append({
            "id": s.id, "srn": s.srn, "name": s.name,
            "status": att["status"] if att else None,
            "marked_at": att["timestamp"].isoformat() if att else None
        })

    return jsonify({"success": True, "students": out})


@app.route("/faculty/attendance", methods=["POST", "PUT"])
@roles_allowed(["professor", "admin"])
def mark_attendance():
    payload = request.json or {}
    fid = int(get_jwt_identity())
    
    if request.method == "POST":
        # Mark Attendance
        items = payload.get("data") # List of {student_id, status}
        degree = payload.get("degree")
        semester = payload.get("semester")
        section = payload.get("section")
        subject = payload.get("subject")
        date_str = payload.get("date") # YYYY-MM-DD
        
        if not (items and degree and semester and section and subject and date_str):
            return jsonify({"success": False, "message": "Missing fields"}), 400
            
        try:
            date_obj = datetime.strptime(date_str, "%Y-%m-%d").date()
        except ValueError:
            return jsonify({"success": False, "message": "Invalid date format"}), 400

        # Check existing logic? Optional: Prevent double marking or just overwrite?
        # User requested "editable for 30 mins".
        
        count = 0
        for item in items:
            sid = item.get("student_id")
            status = item.get("status")
            
            # Check existing
            exist = Attendance.query.filter_by(student_id=sid, subject=subject, date=date_obj).first()
            if exist:
                # Check edit window
                window = datetime.utcnow() - exist.timestamp
                if window.total_seconds() > 1800: # 30 mins
                     # Skip or Error? Let's skip locked ones but continue others?
                     # Ideally should block whole batch if one is locked?
                     # Let's simple check: If ANY locked, abort?
                     # Implementation: Just overwrite if within window.
                     continue 
                else:
                    exist.status = status
                    exist.timestamp = datetime.utcnow() # Reset timer on edit? Or keep original? "editable for a certain amount of time AFTER MARKING" -> implies original mark time.
                    # If I reset timestamp, window extends. Let's NOT update timestamp on edit to enforce strict 30m from FIRST mark.
                    # But if I created it now, I set timestamp.
            else:
                new_att = Attendance(
                    student_id=sid, faculty_id=fid, degree=degree, 
                    semester=int(semester), section=section, subject=subject, 
                    date=date_obj, status=status
                )
                db.session.add(new_att)
            
            # Notify if Absent
            if status == "Absent":
                stu = User.query.get(sid)
                if stu:
                    details = {"Subject": subject, "Date": date_str, "Status": "Absent"}
                    # Notify Student
                    send_professional_email(stu.email, f"Attendance Alert: Absent for {subject}", "Absence Recorded", details, f"You have been marked <strong>Absent</strong> for {subject} on {date_str}.")
                    # Notify Parent
                    if stu.parent_email:
                        send_professional_email(stu.parent_email, f"Attendance Alert: {stu.name} Absent", f"Absence Alert for {stu.name}", details, f"Your ward <strong>{stu.name}</strong> was marked Absent for {subject} on {date_str}.")

            count += 1
            
        db.session.commit()
        return jsonify({"success": True, "message": f"Attendance marked for {count} students."})

    elif request.method == "PUT":
        # Single Edit (if needed) or Bulk Edit same as POST logic?
        # I'll reuse POST for "Marking/Editing" as specified in requirement usually implies one interface.
        return jsonify({"success": False, "message": "Use POST to mark or update."})


@app.route("/student/attendance", methods=["GET"])
@roles_allowed(["student", "parent"])
def get_my_attendance():
    uid = int(get_jwt_identity())
    # Return list of {date, subject, status}
    rows = Attendance.query.filter_by(student_id=uid).order_by(Attendance.date.desc()).all()
    out = [{
        "date": r.date.isoformat(),
        "subject": r.subject,
        "status": r.status
    } for r in rows]
    return jsonify({"success": True, "attendance": out})


# -------------------- AI CHAT (Gemini) --------------------

def call_groq_api(prompt: str):
    # Check if the key is loaded and present
    if not GROQ_API_KEY:
        return None, "GROQ_API_KEY environment variable is missing or empty."
        
    API_URL = "https://api.groq.com/openai/v1/chat/completions"
    MAX_RETRIES = 3
    payload = {
        "model": "llama-3.3-70b-versatile",
        "messages": [
            {"role": "system", "content": "You are Orbit Bot, an Academic Assistant trained by Meta and tuned at LeafCore Labs. If asked who you are, introduce yourself using this identity. Provide helpful, concise, and academically relevant answers."},
            {"role": "user", "content": prompt}
        ]
    }
    for attempt in range(MAX_RETRIES):
        try:
            response = requests.post(API_URL, headers={'Authorization': f'Bearer {GROQ_API_KEY}', 'Content-Type': 'application/json'}, data=json.dumps(payload))
            response.raise_for_status()
            result = response.json()
            text = result.get('choices', [{}])[0].get('message', {}).get('content')
            if text:
                return text, None
            else:
                return None, f"API returned an empty response. Response JSON: {result}"
        except requests.exceptions.HTTPError as e:
            if response.status_code in [429, 500, 503] and attempt < MAX_RETRIES - 1:
                time.sleep(2 ** attempt)
                continue
            return None, f"HTTP Error: {e} - Response Text: {response.text}"
        except requests.exceptions.RequestException as e:
            return None, f"Request failed: {e}"
        except Exception as e:
            return None, f"An unexpected error occurred: {e}"
    return None, f"Failed to connect to API after {MAX_RETRIES} attempts."


@app.route("/chat", methods=["POST"])
@jwt_required()
def chat():
    # Handle both JSON and Multipart
    q = ""
    file_text = ""
    session_id = ""
    user_id = int(get_jwt_identity())
    
    # 1. Parse Input
    if request.is_json:
        data = request.json or {}
        q = data.get("question") or data.get("q") or ""
        session_id = data.get("session_id")
    else:
        q = request.form.get("question") or request.form.get("q") or ""
        session_id = request.form.get("session_id")
        
        # 2. Handle File Upload
        if 'file' in request.files:
            file = request.files['file']
            if file and file.filename:
                filename = file.filename.lower()
                try:
                    if filename.endswith('.pdf'):
                        try:
                            import PyPDF2
                            pdf_reader = PyPDF2.PdfReader(file.stream)
                            for page in pdf_reader.pages:
                                text = page.extract_text()
                                if text:
                                    file_text += text + "\n"
                        except ImportError:
                             return jsonify({"success": False, "message": "Server missing PyPDF2 library. Please install it."}), 500
                    
                    elif filename.endswith('.docx') or filename.endswith('.doc'):
                        if not docx: return jsonify({"success": False, "message": "Server missing python-docx library."}), 500
                        try:
                            doc = docx.Document(file.stream)
                            for para in doc.paragraphs:
                                file_text += para.text + "\n"
                        except Exception as e:
                            file_text = f"[Error reading DOCX: {str(e)}]"

                    elif filename.endswith('.pptx'):
                        if not Presentation: return jsonify({"success": False, "message": "Server missing python-pptx library."}), 500
                        try:
                            prs = Presentation(file.stream)
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        file_text += shape.text + "\n"
                        except Exception as e:
                             file_text = f"[Error reading PPTX: {str(e)}]"



                    elif filename.endswith(('.txt', '.md', '.py', '.js', '.html', '.css', '.json')):
                         file_text = file.read().decode('utf-8', errors='ignore')
                    else:
                        file_text = f"[Uploaded file: {file.filename} - Type not supported for automatic reading]"
                except Exception as e:
                    print(f"File read error: {e}")
                    file_text = f"[Error reading file: {str(e)}]"

    if not q and not file_text:
        return jsonify({"success": False, "message": "Question or file required"}), 400

    # 3. Session Management
    session = None
    if session_id:
        session = AIChatSession.query.filter_by(id=session_id, user_id=user_id).first()
        if not session:
             return jsonify({"success": False, "message": "Session not found or access denied"}), 404
        session.updated_at = datetime.utcnow()
    else:
        # Create New Session
        # Generate title from query (first 50 chars) or "New Chat" if empty/file-only
        title = q[:50] + "..." if q else "File Analysis"
        session = AIChatSession(user_id=user_id, title=title)
        db.session.add(session)
        db.session.commit()
    
    # 4. Save User Message
    user_msg_text = q
    if file_text:
        user_msg_text += f"\n[Attached: {request.files['file'].filename if 'file' in request.files else 'File'}]"
    
    db.session.add(AIChatMessage(session_id=session.id, role="user", text=user_msg_text))
    
    # 5. Construct Prompt with Context
    final_prompt = q
    if file_text:
        final_prompt += f"\n\n--- CONTEXT FROM UPLOADED FILE ({request.files['file'].filename}) ---\n{file_text[:20000]}\n--- END CONTEXT ---\n(Note: Text has been extracted from the file. Answer based on this context if relevant.)"

    # Context Awareness: Retrieve last few messages for context?
    # For now, keeping it stateless per request to save tokens, but could fetch history here.
    
    answer, error_msg = call_groq_api(final_prompt)
    
    if answer:
        # Save AI Response
        db.session.add(AIChatMessage(session_id=session.id, role="ai", text=answer))
        db.session.commit()
        return jsonify({"success": True, "answer": answer, "session_id": session.id})
    else:
        return jsonify({"success": False, "message": error_msg or "Failed to get response from AI model."}), 500


@app.route("/ai/sessions", methods=["GET"])
@jwt_required()
def get_ai_sessions():
    uid = int(get_jwt_identity())
    sessions = AIChatSession.query.filter_by(user_id=uid).order_by(AIChatSession.updated_at.desc()).all()
    out = [{"id": s.id, "title": s.title, "updated_at": s.updated_at.isoformat()} for s in sessions]
    return jsonify({"success": True, "sessions": out})


@app.route("/ai/session/<session_id>", methods=["GET", "DELETE"])
@jwt_required()
def manage_ai_session(session_id):
    uid = int(get_jwt_identity())
    session = AIChatSession.query.filter_by(id=session_id, user_id=uid).first()
    if not session:
        return jsonify({"success": False, "message": "Session not found"}), 404

    if request.method == 'DELETE':
        try:
            AIChatMessage.query.filter_by(session_id=session_id).delete() # Manual cascade
            db.session.delete(session)
            db.session.commit()
            return jsonify({"success": True, "message": "Session deleted"})
        except Exception as e:
            print(f"Error deleting session: {e}")
            db.session.rollback()
            return jsonify({"success": False, "message": f"Server error: {str(e)}"}), 500

    messages = AIChatMessage.query.filter_by(session_id=session_id).order_by(AIChatMessage.timestamp.asc()).all()
    out = [{"role": m.role, "text": m.text, "timestamp": m.timestamp.isoformat()} for m in messages]
    return jsonify({"success": True, "messages": out, "title": session.title})


@app.route("/api/academic-insights", methods=["GET"])
@jwt_required()
def get_academic_insights():
    """
    Aggregates student data (Marks, Attendance, Feedback) and asks Groq AI for insights.
    Returns: JSON with risks, priorities, and persona-based advice.
    """
    current_uid = int(get_jwt_identity())
    user = User.query.get(current_uid)
    
    # If parent, get ward's ID
    student_id = current_uid
    if user.role == "parent":
        # Parent login sends SRN as 'email' usually, but here we need to find the student linked to this parent.
        # Current login logic returns PARENT user object (which is a student object but role=parent? No, we have separate parent login flow).
        # Wait, the parent login logic in /login issues token with identity=student.id and role='parent'.
        # So get_jwt_identity() is ALREADY the student_id.
        pass
    
    # 1. Fetch Marks
    marks = Mark.query.filter_by(student_id=student_id).all()
    marks_summary = [{"subject": m.subject, "score": m.marks_obtained, "max": m.max_marks} for m in marks]
    
    # 2. Fetch Attendance
    # Calculate % per subject
    attendance_records = Attendance.query.filter_by(student_id=student_id).all()
    att_stats = {} # {subject: {present: 0, total: 0}}
    for r in attendance_records:
        if r.subject not in att_stats: att_stats[r.subject] = {"present": 0, "total": 0}
        att_stats[r.subject]["total"] += 1
        if r.status == "Present":
            att_stats[r.subject]["present"] += 1
            
    att_summary = []
    for subj, data in att_stats.items():
        pct = (data["present"] / data["total"]) * 100 if data["total"] > 0 else 0
        att_summary.append({"subject": subj, "percentage": round(pct, 1)})
        
    # 3. Fetch Feedback
    feedbacks = Feedback.query.filter_by(student_id=student_id).order_by(Feedback.created_at.desc()).limit(3).all()
    fb_summary = [{"subject": f.subject, "text": f.text} for f in feedbacks]
    
    # 4. Construct Prompt
    role_view = "Student" if user.role == "student" else "Parent"
    prompt = f"""
    Analyze the academic data for a student named {user.name}.
    Role View: {role_view} (Provide advice suitable for a {role_view}).
    
    Data:
    Marks: {json.dumps(marks_summary)}
    Attendance: {json.dumps(att_summary)}
    Recent Feedback: {json.dumps(fb_summary)}
    
    Task:
    1. Identify 'Attendance Risks' (Subjects < 75%).
    2. Identify 'Subject Priorities' (Low marks).
    3. Generate 'Improvement Suggestions' (Actionable steps).
    4. Write a 'Counselor Message': A warm, professional paragraph summarizing status and advice.
       - If Student view: Be encouraging, strategic.
       - If Parent view: Be reassuring, clear on what to monitor.
       
    Output strictly in valid JSON format:
    {{
        "attendance_risks": ["Subject A", ...],
        "priorities": ["Subject B", ...],
        "suggestions": ["Tip 1", "Tip 2", ...],
        "counselor_message": "..."
    }}
    """
    
    # 5. Call AI
    ai_response, error = call_groq_api(prompt)
    
    if not ai_response:
        # Fallback if AI fails
        return jsonify({
            "success": True, 
            "insights": {
                "attendance_risks": [], "priorities": [], 
                "suggestions": ["Focus on consistent attendance.", "Review recent class notes."],
                "counselor_message": "AI services are currently unavailable, but please review your marks and attendance manually."
            }
        })
        
    # Parse JSON from AI (It might wrap in markdown code blocks)
    try:
        clean_json = ai_response.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(clean_json)
        return jsonify({"success": True, "insights": parsed})
    except Exception as e:
        print(f"AI Parse Error: {e}")
        return jsonify({"success": True, "insights": {
             "attendance_risks": [], "priorities": [], 
             "suggestions": ["Please review your dashboard."], 
             "counselor_message": ai_response # Return raw text if JSON parse fails
        }})


# -------------------- DB INIT --------------------

def init_db():
    db.create_all() # This now includes Hostel, Room, and HostelAllocation
    admin_email = DEFAULT_ADMIN_EMAIL
    admin_password = DEFAULT_ADMIN_PASSWORD
    if not User.query.filter_by(email=admin_email).first():
        u = User(
            srn=None, name="Default Admin", email=admin_email,
            password_hash=hash_password(admin_password), role="admin", status="APPROVED"
        )
        db.session.add(u); db.session.commit(); print("Created default admin:", admin_email)
    if Degree.query.count() == 0:
        for d in ["BCA", "BE", "MCA", "MBA"]:
            db.session.add(Degree(name=d))
        db.session.commit()
    if Section.query.count() == 0:
        for s in ["A", "B", "C"]:
            db.session.add(Section(name=s))
        db.session.commit()
        
    # NEW: Initial Hostel/Room for testing (Optional but helpful)
    if Hostel.query.count() == 0:
        h1 = Hostel(name="Boys Hostel A", address="South Campus")
        h2 = Hostel(name="Girls Hostel B", address="North Campus")
        db.session.add_all([h1, h2])
        db.session.flush()
        if h1.id:
            db.session.add(Room(hostel_id=h1.id, room_number="A101", capacity=2))
            db.session.add(Room(hostel_id=h1.id, room_number="A102", capacity=3))
        if h2.id:
            db.session.add(Room(hostel_id=h2.id, room_number="B201", capacity=2))
        db.session.commit()
        print("Created default hostels and rooms.")

    # Create Messages table if it doesn't exist (handled by create_all normally, but just to be sure if adding to existing DB)
    # db.create_all() handles it
    
    # HRD (Placement Cell) Default Admin Creation
    hrd_email = HRD_EMAIL
    hrd_password = HRD_PASSWORD
    if not HRDUser.query.filter_by(email=hrd_email).first():
        hrd_admin = HRDUser(
            email=hrd_email,
            password_hash=hash_password(hrd_password),
            name="HRD Admin",
            phone=None,
            department="Placement Cell",
            is_active=True
        )
        db.session.add(hrd_admin)
        db.session.commit()
        print(f"Created default HRD admin: {hrd_email}")

    os.makedirs(RECEIPT_TMP_DIR, exist_ok=True)


# -------------------- PARENT COMMUNICATION ROUTES --------------------

@app.route("/parent/professors", methods=["GET"])
@roles_allowed(["parent", "student"]) # Allow both to see their professors
def parent_get_professors():
    identity = get_jwt_identity()
    user = User.query.get(identity)
    if not user:
        return jsonify({"success": False, "message": "User not found"}), 404

    # Find allocations for the student's current class
    allocations = FacultyAllocation.query.filter_by(
        degree=user.degree,
        semester=user.semester,
        section=user.section
    ).all()

    professors = []
    seen_ids = set()
    
    for alloc in allocations:
        prof = User.query.get(alloc.faculty_id)
        if prof and prof.id not in seen_ids:
            professors.append({
                "id": prof.id,
                "name": prof.name,
                "email": prof.email,
                "subject": alloc.subject, # Primary subject for context
                "allocations": [a.subject for a in allocations if a.faculty_id == prof.id] # All subjects taught by this prof to this class
            })
            seen_ids.add(prof.id)

    return jsonify({"success": True, "professors": professors})


@app.route("/parent/contact-professor", methods=["POST"])
@roles_allowed(["parent"]) 
def parent_contact_professor():
    identity = get_jwt_identity()
    student = User.query.get(identity) # Parent token identity is the Student ID
    if not student:
        return jsonify({"success": False, "message": "Student record not found"}), 404

    data = request.json or {}
    prof_id = data.get("professor_id")
    subject_line = data.get("subject", "Parent Inquiry")
    message_body = data.get("message", "")
    
    if not prof_id or not message_body:
        return jsonify({"success": False, "message": "Professor and message are required"}), 400

    prof = User.query.get(prof_id)
    if not prof:
         return jsonify({"success": False, "message": "Professor not found"}), 404

    # Construct the Email
    # We use the Parent Email from the student record if available, else fallback
    reply_to = student.parent_email or f"parent_of_{student.srn}@noteorbit.com" 
    
    email_subject = f"[NoteOrbit] Parent Query: {student.name} ({student.srn})"
    
    html_content = f"""
    <div style="font-family: Arial, sans-serif; color: #333; max-width: 600px; margin: 0 auto; border: 1px solid #eee; border-radius: 8px; overflow: hidden;">
        <div style="background-color: #3b82f6; padding: 20px; color: white;">
            <h2 style="margin: 0;">Parent Communication</h2>
            <p style="margin: 5px 0 0; opacity: 0.9;">From the desk of {student.name}'s Guardian</p>
        </div>
        <div style="padding: 30px;">
            <p><strong>To:</strong> Prof. {prof.name}</p>
            <p><strong>Regarding Student:</strong> {student.name} ({student.srn})</p>
            <p><strong>Class:</strong> {student.degree} - Sem {student.semester} (Sec {student.section})</p>
            <hr style="border: 0; border-top: 1px solid #eee; margin: 20px 0;">
            
            <h3 style="color: #3b82f6; margin-top: 0;">{subject_line}</h3>
            <p style="background-color: #f9fafb; padding: 15px; border-radius: 6px; border-left: 4px solid #3b82f6;">
                "{message_body}"
            </p>
            
            <p style="margin-top: 30px; font-size: 13px; color: #666;">
                You can reply directly to this email to contact the parent at <a href="mailto:{reply_to}">{reply_to}</a>.
            </p>
        </div>
         <div style="background-color: #f4f4f9; padding: 15px; text-align: center; font-size: 12px; color: #999;">
            NoteOrbit Academic Portal
        </div>
    </div>
    """

    # We can't strictly set "Reply-To" in our simple helper without modification, 
    # but we can try injecting it into headers if we modify send_email later. 
    # For now, we rely on the body containing the contact info.
    # Actually, let's just assume the Prof sees the "From" address as the System and replies to the text in body.
    
    # --- SAVE TO DB (In-App Messaging) ---
    try:
        new_msg = Message(
            sender="parent",
            student_id=student.id,
            faculty_id=prof.id,
            subject=subject_line,
            body=message_body
        )
        db.session.add(new_msg)
        db.session.commit()
    except Exception as e:
        print(f"Error saving message to DB: {e}")
    
    # Send Email Notification
    send_email(prof.email, email_subject, html_content)
    
    return jsonify({"success": True, "message": f"Message sent to Prof. {prof.name}"})


@app.route("/faculty/conversations", methods=["GET"])
@roles_allowed(["professor"])
def get_faculty_conversations():
    """Returns list of students the faculty has chatted with, ordered by most recent activity."""
    identity = get_jwt_identity()
    prof_id = int(identity)
    
    # Get all messages for this faculty
    msgs = Message.query.filter_by(faculty_id=prof_id).order_by(Message.created_at.desc()).all()
    
    # Group by student
    students_map = {}
    for m in msgs:
        if m.student_id not in students_map:
            s_obj = User.query.get(m.student_id)
            if not s_obj: continue
            
            # Count unread from parent
            unread_count = 0
            
            students_map[m.student_id] = {
                "student_id": s_obj.id,
                "student_name": s_obj.name,
                "student_srn": s_obj.srn,
                "last_message": m.body[:50] + "..." if len(m.body)>50 else m.body,
                "last_timestamp": m.created_at.isoformat(),
                "unread_count": 0
            }
        
        # Increment unread if message is from parent and not read
        if m.sender == 'parent' and not m.is_read:
            students_map[m.student_id]["unread_count"] += 1

    return jsonify({"success": True, "conversations": list(students_map.values())})


@app.route("/faculty/messages/<int:student_id>", methods=["GET"])
@roles_allowed(["professor"])
def get_conversation_thread(student_id):
    """Gets full chat history with a specific student/parent."""
    identity = get_jwt_identity()
    prof_id = int(identity)

    msgs = Message.query.filter_by(faculty_id=prof_id, student_id=student_id).order_by(Message.created_at.asc()).all()
    
    # Mark all parent messages as read
    for m in msgs:
        if m.sender == 'parent' and not m.is_read:
            m.is_read = True
    db.session.commit()
    
    out = []
    for m in msgs:
        out.append({
            "id": m.id,
            "sender": m.sender, # 'parent' or 'faculty'
            "body": m.body,
            "subject": m.subject,
            "timestamp": m.created_at.isoformat()
        })
        
    return jsonify({"success": True, "messages": out})


@app.route("/faculty/messages/reply", methods=["POST"])
@roles_allowed(["professor"])
def reply_to_parent():
    identity = get_jwt_identity()
    prof = User.query.get(identity)
    
    data = request.json or {}
    student_id = data.get("student_id") # Use Student ID context
    reply_body = data.get("reply_body")
    
    if not student_id or not reply_body:
        return jsonify({"success": False, "message": "Student ID and reply body required"}), 400
        
    student = User.query.get(student_id)
    if not student or not student.parent_email:
        return jsonify({"success": False, "message": "Parent email not found"}), 404
        
    # 1. Save to DB
    try:
        new_msg = Message(
            sender="faculty",
            student_id=student.id,
            faculty_id=prof.id,
            subject=f"Reply: Parent Query ({student.name})",
            body=reply_body,
            is_read=True # Faculty read their own message
        )
        db.session.add(new_msg)
        db.session.commit()
    except Exception as e:
        print("DB Save Error:", e)

    # 2. Send Email
    subject = f"Re: Query regarding {student.name} - [Prof. {prof.name}]"
    html_content = f"""
    <div style="font-family: Arial, sans-serif; padding: 20px;">
        <p><strong>From:</strong> Prof. {prof.name}</p>
        <hr/>
        <p>{reply_body}</p>
        <p style="color: #666; font-size: 12px; margin-top: 20px;">
           Reply strictly via email or check the portal.
        </p>
    </div>
    """
    send_email(student.parent_email, subject, html_content)
    return jsonify({"success": True, "message": "Reply sent"})


@app.route("/parent/conversations", methods=["GET"])
@roles_allowed(["parent", "student"]) 
def get_parent_conversations():
    """Returns list of professors the student/parent has chatted with."""
    identity = get_jwt_identity()
    student = User.query.get(identity)
    
    # Get all messages where this student is involved
    msgs = Message.query.filter_by(student_id=student.id).order_by(Message.created_at.desc()).all()
    
    faculty_map = {}
    for m in msgs:
        if m.faculty_id not in faculty_map:
            prof = User.query.get(m.faculty_id)
            if not prof: continue
            
            # Count unread from faculty
            unread_count = 0
            
            faculty_map[m.faculty_id] = {
                "faculty_id": prof.id,
                "faculty_name": prof.name,
                "last_message": m.body[:50] + "..." if len(m.body)>50 else m.body,
                "last_timestamp": m.created_at.isoformat(),
                "unread_count": 0
            }
        
        if m.sender == 'faculty' and not m.is_read:
            faculty_map[m.faculty_id]["unread_count"] += 1

    return jsonify({"success": True, "conversations": list(faculty_map.values())})


@app.route("/parent/messages/<int:faculty_id>", methods=["GET"])
@roles_allowed(["parent", "student"])
def get_parent_thread(faculty_id):
    """Gets full chat history with a specific professor."""
    identity = get_jwt_identity()
    student_id = int(identity)

    msgs = Message.query.filter_by(student_id=student_id, faculty_id=faculty_id).order_by(Message.created_at.asc()).all()
    
    # Mark all faculty messages as read
    for m in msgs:
        if m.sender == 'faculty' and not m.is_read:
            m.is_read = True
    db.session.commit()
    
    out = []
    for m in msgs:
        out.append({
            "id": m.id,
            "sender": m.sender, # 'parent' or 'faculty'
            "body": m.body,
            "subject": m.subject,
            "timestamp": m.created_at.isoformat()
        })
        
    return jsonify({"success": True, "messages": out})


@app.route("/parent/messages/reply", methods=["POST"])
@roles_allowed(["parent"]) # Only parent can reply, student read-only? Let's allow parent.
def parent_reply():
    identity = get_jwt_identity()
    student = User.query.get(identity)
    
    data = request.json or {}
    faculty_id = data.get("faculty_id")
    reply_body = data.get("reply_body")
    
    if not faculty_id or not reply_body:
        return jsonify({"success": False, "message": "Faculty ID and reply body required"}), 400
        
    prof = User.query.get(faculty_id)
    if not prof:
        return jsonify({"success": False, "message": "Professor not found"}), 404

    # 1. Save to DB
    try:
        new_msg = Message(
            sender="parent",
            student_id=student.id,
            faculty_id=prof.id,
            subject=f"Reply from Parent of {student.name}",
            body=reply_body,
            is_read=False
        )
        db.session.add(new_msg)
        db.session.commit()
    except Exception as e:
        print("DB Save Error:", e)

    # 2. Send Email to Faculty
    subject = f"New Message from Parent of {student.name}"
    html_content = f"""
    <div style="padding: 20px;">
        <p><strong>Parent Reply:</strong></p>
        <p>{reply_body}</p>
        <p style="color: #666; font-size: 12px; margin-top: 20px;">Login to portal to reply.</p>
    </div>
    """
    send_email(prof.email, subject, html_content)
    
    return jsonify({"success": True, "message": "Reply sent"})


# -------------------- START --------------------


# -------------------- ATTENDANCE FEATURE --------------------
attendance_bp = Blueprint('attendance', __name__)

def create_daily_message_with_ai(msg_type, user_name):
    """
    Generates a fun message using Groq AI.
    msg_type: 'sunday' | 'no_class'
    """
    if not GROQ_API_KEY:
        return "Enjoy your day! (AI Key missing)"

    prompt = ""
    if msg_type == 'sunday':
        prompt = f"Write a short, relaxing, and funny message for a student named {user_name} because it's Sunday. Max 2 sentences. Include a relaxing emoji."
    elif msg_type == 'no_class':
        prompt = f"Write a short, hype message for a student named {user_name} who has no classes today. Include a 'Did you know?' fun fact. Max 3 sentences."

    try:
        resp = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {GROQ_API_KEY}",
                "Content-Type": "application/json"
            },
            json={
                "model": "llama3-8b-8192",
                "messages": [{"role": "user", "content": prompt}]
            },
            timeout=5
        )
        if resp.status_code == 200:
            return resp.json()['choices'][0]['message']['content']
    except Exception as e:
        print(f"AI Gen Error: {e}")
    
    return "Have a great day!"

@attendance_bp.route('/routine/upload', methods=['POST'])
@jwt_required()
def upload_routine():
    user_id = get_jwt_identity()
    user = User.query.get(user_id)
    
    # 1. Get Content (Text or File)
    raw_text = request.form.get('routine_text', '')
    file = request.files.get('file')

    if file:
        try:
            filename = secure_filename(file.filename).lower()
            if filename.endswith('.txt'):
                raw_text = file.read().decode('utf-8')
            elif filename.endswith('.pdf'):
                try:
                    import pypdf
                    pdf = pypdf.PdfReader(file)
                    raw_text = " ".join([page.extract_text() for page in pdf.pages])
                except ImportError:
                    return jsonify({"success": False, "message": "Server missing 'pypdf' library"}), 500
            elif filename.endswith(('.png', '.jpg', '.jpeg', '.webp')):
                # Use Vision Model for Images
                import base64
                image_data = base64.b64encode(file.read()).decode('utf-8')
                
                try:
                    vision_resp = requests.post(
                        "https://api.groq.com/openai/v1/chat/completions",
                        headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
                        json={
                            "model": "llama-3.2-11b-vision-preview",
                            "messages": [
                                {
                                    "role": "user",
                                    "content": [
                                        {"type": "text", "text": "Analyze this image of a class schedule. Extract every single piece of text visible, maintaining the structure as much as possible to help identify which classes belong to which days and times. Output only the raw text content."},
                                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{image_data}"}}
                                    ]
                                }
                            ]
                        }
                    )
                    if vision_resp.status_code != 200:
                        raise Exception(f"Groq Vision API Error ({vision_resp.status_code}): {vision_resp.text}")
                        
                    raw_text = vision_resp.json()['choices'][0]['message']['content']
                except Exception as e:
                    return jsonify({"success": False, "message": f"Vision AI Failed: {str(e)}"}), 500
        except Exception as e:
            return jsonify({"success": False, "message": f"File processing failed: {str(e)}"}), 400

    if not raw_text:
        return jsonify({"success": False, "message": "No routine content found"}), 400

    # 2. AI Parsing (Standard Logic)
    parsed_routine = {}
    try:
        sys_prompt = """You are a strict JSON data extraction assistant. Your task is to extract a weekly class routine from the provided unstructured text (which may contain OCR errors).
    1. Identify days of the week (Monday, Tuesday, etc.).
    2. For each day, list the classes in chronological order.
    3. Format each class as simple text: 'SubjectName (Time)'. Example: 'Mathematics (10:00 AM)'.
    4. Return ONLY a valid JSON object. Keys must be full English day names (e.g., 'Monday'). Values must be Lists of strings.
    5. Ignore unrelated text, headers, or footers.
    6. If a day has no classes, do not include it in the JSON.
    7. Do not wrap the output in markdown code blocks.
    8. Use standard full day names: Monday, Tuesday, Wednesday, Thursday, Friday, Saturday, Sunday.
        """
        
        resp = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
            json={
                "model": "llama3-70b-8192", # Stronger model for logic
                "messages": [
                    {"role": "system", "content": sys_prompt},
                    {"role": "user", "content": f"Extract routine from this text:\n\n{raw_text}"}
                ],
                "response_format": {"type": "json_object"}
            }
        )
        if resp.status_code != 200:
            raise Exception(f"Groq Parsing API Error ({resp.status_code}): {resp.text}")

        content = resp.json()['choices'][0]['message']['content']
        parsed_routine = json.loads(content)
    except Exception as e:
        return jsonify({"success": False, "message": f"AI Parsing Failed: {str(e)}"}), 500

    # 3. Save/Update Routine
    try:
        # Clear old routine
        StudentRoutine.query.filter_by(user_id=user.id).delete()
        
        for day, subs in parsed_routine.items():
            if isinstance(subs, list) and subs:
                new_r = StudentRoutine(
                    user_id=user.id,
                    day_of_week=day,
                    subjects=json.dumps(subs)
                )
                db.session.add(new_r)
        
        db.session.commit()
        return jsonify({"success": True, "message": "Routine updated from file/text", "routine": parsed_routine})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@attendance_bp.route('/routine', methods=['GET'])
@jwt_required()
def get_routine():
    user_id = get_jwt_identity()
    routines = StudentRoutine.query.filter_by(user_id=user_id).all()
    out = {}
    for r in routines:
        out[r.day_of_week] = json.loads(r.subjects)
    return jsonify({"success": True, "routine": out})

@attendance_bp.route('/routine', methods=['DELETE'])
@jwt_required()
def delete_routine():
    user_id = get_jwt_identity()
    StudentRoutine.query.filter_by(user_id=user_id).delete()
    db.session.commit()
    return jsonify({"success": True, "message": "Routine deleted"})

@attendance_bp.route('/today', methods=['GET'])
@jwt_required()
def get_today_status():
    user_id = get_jwt_identity()
    user = User.query.get(user_id)
    today = datetime.utcnow().date()
    day_name = today.strftime("%A")

    # 1. Check if Sunday
    if day_name == "Sunday":
        msg = create_daily_message_with_ai('sunday', user.name)
        return jsonify({
            "status": "holiday",
            "message": msg,
            "can_mark": False
        })

    # 2. Check if already marked (ANY entry for today)
    logs = StudentAttendanceLog.query.filter_by(user_id=user_id, date=today).all()
    
    if logs:
        # Check if it was "No Class"
        if len(logs) == 1 and logs[0].status == 'No Class':
             msg = create_daily_message_with_ai('no_class', user.name)
             return jsonify({"status": "marked_no_class", "can_mark": False, "fun_message": msg})
        
        # Details
        log_data = [{"subject": l.subject, "status": l.status} for l in logs]
        
        # Determine mood based on attendance
        present_count = sum(1 for l in logs if l.status == 'Present')
        total = len(logs)
        mood = 'good' if (total > 0 and (present_count/total) > 0.75) else 'bad'
        msg = create_daily_message_with_ai(mood, user.name)
        
        return jsonify({"status": "marked", "logs": log_data, "can_mark": False, "fun_message": msg})

    # 3. Not marked yet - Get Routine
    routine = StudentRoutine.query.filter_by(user_id=user.id, day_of_week=day_name).first()
    
    if not routine:
        return jsonify({
            "status": "pending",
            "subjects": [],
            "message": "No routine found for today. Add one first?",
            "can_mark": False
        })

    subjects = json.loads(routine.subjects)
    return jsonify({
        "status": "pending",
        "subjects": subjects,
        "can_mark": True
    })

@attendance_bp.route('/mark', methods=['POST'])
@jwt_required()
def mark_attendance():
    """
    Body: 
    { "type": "classes", "data": {"Math": "Present", "CS": "Absent"} }
    OR
    { "type": "no_class" }
    """
    user_id = get_jwt_identity()
    user = User.query.get(user_id)
    today = datetime.utcnow().date()
    
    # Prevent double marking
    if StudentAttendanceLog.query.filter_by(user_id=user_id, date=today).first():
        return jsonify({"success": False, "message": "Already marked for today"}), 400

    data = request.json
    mark_type = data.get('type')

    try:
        if mark_type == 'no_class':
            # Log single entry
            log = StudentAttendanceLog(
                user_id=user.id,
                date=today,
                subject=None,
                status='No Class'
            )
            db.session.add(log)
            msg = create_daily_message_with_ai('no_class', user.name)
            db.session.commit()
            return jsonify({"success": True, "message": "Enjoy your day!", "fun_message": msg})

        elif mark_type == 'classes':
            attendance_map = data.get('data', {})
            present_count = 0
            total_count = 0
            
            for sub, status in attendance_map.items():
                log = StudentAttendanceLog(
                    user_id=user.id,
                    date=today,
                    subject=sub,
                    status=status
                )
                db.session.add(log)
                if status == 'Present': present_count += 1
                total_count += 1
                
            db.session.commit()
            
            # Generate fun message
            mood = 'good' if (total_count > 0 and (present_count/total_count) > 0.75) else 'bad'
            msg = create_daily_message_with_ai(mood, user.name)
            
            return jsonify({"success": True, "message": "Attendance Saved", "fun_message": msg})
        
            
            return jsonify({"success": True, "message": "Attendance Saved", "fun_message": msg})
        
        else:
            return jsonify({"success": False, "message": "Invalid type"}), 400

    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@attendance_bp.route('/today', methods=['DELETE'])
@jwt_required()
def reset_today():
    user_id = get_jwt_identity()
    today = datetime.utcnow().date()
    
    try:
        StudentAttendanceLog.query.filter_by(user_id=user_id, date=today).delete()
        db.session.commit()
        return jsonify({"success": True, "message": "Today's attendance reset."})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@attendance_bp.route('/stats', methods=['GET'])
@jwt_required()
def get_stats():
    user_id = get_jwt_identity()
    
    # Fetch all logs except 'No Class' or 'Holiday'
    logs = StudentAttendanceLog.query.filter(
        StudentAttendanceLog.user_id == user_id,
        StudentAttendanceLog.status.in_(['Present', 'Absent'])
    ).all()

    total_classes = len(logs)
    if total_classes == 0:
        return jsonify({"success": True, "overall": 0, "subject_wise": []})

    present_count = sum(1 for l in logs if l.status == 'Present')
    overall = (present_count / total_classes) * 100

    # Subject-wise
    sub_stats = {}
    for l in logs:
        if l.subject not in sub_stats:
            sub_stats[l.subject] = {"present": 0, "total": 0}
        sub_stats[l.subject]["total"] += 1
        if l.status == 'Present':
            sub_stats[l.subject]["present"] += 1
    
    # Calculate %
    final_subs = []
    for sub, dat in sub_stats.items():
        pct = (dat["present"] / dat["total"]) * 100
        final_subs.append({
            "subject": sub,
            "present": dat["present"],
            "total": dat["total"],
            "percentage": round(pct, 1)
        })

    return jsonify({
        "success": True,
        "overall": round(overall, 1),
        "subject_wise": final_subs,
        "history": [] # TODO: Add history list if needed
    })

# Register Blueprint
# ==================== HRD (PLACEMENT CELL) REST API ROUTES ====================

# ----- HRD AUTHENTICATION -----
@app.route("/hrd/login", methods=["POST"])
def hrd_login():
    """HRD login endpoint"""
    data = request.json or {}
    email = data.get("email")
    password = data.get("password")
    
    if not email or not password:
        return jsonify({"success": False, "message": "Email and password required"}), 400
    
    hrd_user = HRDUser.query.filter_by(email=email).first()
    if not hrd_user or hrd_user.password_hash != hash_password(password):
        return jsonify({"success": False, "message": "Invalid credentials"}), 401
    
    if not hrd_user.is_active:
        return jsonify({"success": False, "message": "Account inactive"}), 403
    
    token = create_access_token(identity=str(hrd_user.id), additional_claims={"role": "hrd"})
    
    return jsonify({
        "success": True,
        "token": token,
        "user": {
            "id": hrd_user.id,
            "name": hrd_user.name,
            "email": hrd_user.email,
            "department": hrd_user.department,
            "role": "hrd"
        }
    })

# ----- COMPANY MANAGEMENT -----
@app.route("/hrd/companies", methods=["POST"])
@hrd_required
def create_company():
    """Create a new company"""
    data = request.json or {}
    
    try:
        company = Company(
            name=data.get("name"),
            sector=data.get("sector"),
            website=data.get("website"),
            hr_name=data.get("hr_name"),
            hr_email=data.get("hr_email"),
            hr_phone=data.get("hr_phone"),
            visit_history=data.get("visit_history", [])
        )
        db.session.add(company)
        db.session.commit()
        
        hrd_id = get_jwt_identity()
        log_hrd_activity(hrd_id, "hrd", "created_company", "company", company.id, {"name": company.name})
        
        return jsonify({"success": True, "company_id": company.id})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/hrd/companies", methods=["GET"])
@hrd_required
def get_companies():
    """Get all companies"""
    companies = Company.query.filter_by(is_active=True).all()
    return jsonify({
        "success": True,
        "companies": [{
            "id": c.id,
            "name": c.name,
            "sector": c.sector,
            "website": c.website,
            "hr_name": c.hr_name,
            "hr_email": c.hr_email,
            "hr_phone": c.hr_phone,
            "visit_history": c.visit_history,
            "created_at": c.created_at.isoformat()
        } for c in companies]
    })

@app.route("/hrd/companies/<int:company_id>", methods=["GET"])
@hrd_required
def get_company(company_id):
    """Get company details"""
    company = Company.query.get(company_id)
    if not company:
        return jsonify({"success": False, "message": "Company not found"}), 404
    
    return jsonify({
        "success": True,
        "company": {
            "id": company.id,
            "name": company.name,
            "sector": company.sector,
            "website": company.website,
            "hr_name": company.hr_name,
            "hr_email": company.hr_email,
            "hr_phone": company.hr_phone,
            "visit_history": company.visit_history,
            "created_at": company.created_at.isoformat()
        }
    })

@app.route("/hrd/companies/<int:company_id>", methods=["PUT"])
@hrd_required
def update_company(company_id):
    """Update company details"""
    company = Company.query.get(company_id)
    if not company:
        return jsonify({"success": False, "message": "Company not found"}), 404
    
    data = request.json or {}
    
    try:
        company.name = data.get("name", company.name)
        company.sector = data.get("sector", company.sector)
        company.website = data.get("website", company.website)
        company.hr_name = data.get("hr_name", company.hr_name)
        company.hr_email = data.get("hr_email", company.hr_email)
        company.hr_phone = data.get("hr_phone", company.hr_phone)
        company.visit_history = data.get("visit_history", company.visit_history)
        
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/hrd/companies/<int:company_id>", methods=["DELETE"])
@hrd_required
def deactivate_company(company_id):
    """Deactivate company"""
    company = Company.query.get(company_id)
    if not company:
        return jsonify({"success": False, "message": "Company not found"}), 404
    
    company.is_active = False
    db.session.commit()
    return jsonify({"success": True})

# ----- PLACEMENT DRIVE MANAGEMENT -----
@app.route("/hrd/drives", methods=["POST"])
@hrd_required
def create_drive():
    """Create placement drive"""
    data = request.json or {}
    hrd_id = get_jwt_identity()
    
    try:
        drive = PlacementDrive(
            company_id=data.get("company_id"),
            title=data.get("title"),
            description=data.get("description"),
            role=data.get("role"),
            ctc_min=data.get("ctc_min"),
            ctc_max=data.get("ctc_max"),
            location=data.get("location"),
            drive_type=data.get("drive_type", "on_campus"),
            eligibility_criteria=data.get("eligibility_criteria", {}),
            application_start=datetime.fromisoformat(data.get("application_start")) if data.get("application_start") else None,
            application_end=datetime.fromisoformat(data.get("application_end")) if data.get("application_end") else None,
            created_by=hrd_id
        )
        db.session.add(drive)
        db.session.commit()
        
        log_hrd_activity(hrd_id, "hrd", "created_drive", "drive", drive.id, {"title": drive.title})
        
        return jsonify({"success": True, "drive_id": drive.id})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/hrd/drives", methods=["GET"])
@hrd_required
def get_drives():
    """Get all placement drives"""
    drives = PlacementDrive.query.all()
    
    result = []
    for d in drives:
        company = Company.query.get(d.company_id)
        applications_count = DriveApplication.query.filter_by(drive_id=d.id).count()
        
        result.append({
            "id": d.id,
            "company_name": company.name if company else "Unknown",
            "title": d.title,
            "role": d.role,
            "ctc_min": d.ctc_min,
            "ctc_max": d.ctc_max,
            "status": d.status,
            "applications_count": applications_count,
            "application_end": d.application_end.isoformat() if d.application_end else None,
            "created_at": d.created_at.isoformat()
        })
    
    return jsonify({"success": True, "drives": result})

@app.route("/hrd/drives/<int:drive_id>", methods=["GET"])
@hrd_required
def get_drive_details(drive_id):
    """Get drive details with applicants"""
    drive = PlacementDrive.query.get(drive_id)
    if not drive:
        return jsonify({"success": False, "message": "Drive not found"}), 404
    
    company = Company.query.get(drive.company_id)
    applications = DriveApplication.query.filter_by(drive_id=drive_id).all()
    
    applicants = []
    for app in applications:
        student = User.query.get(app.student_id)
        profile = StudentPlacementProfile.query.filter_by(student_id=app.student_id).first()
        
        applicants.append({
            "application_id": app.id,
            "student_id": app.student_id,
            "student_name": student.name if student else "Unknown",
            "student_email": student.email if student else "",
            "degree": student.degree if student else "",
            "semester": student.semester if student else "",
            "status": app.status,
            "applied_at": app.applied_at.isoformat(),
            "tags": app.tags,
            "skills": profile.skills if profile else [],
            "resume_url": profile.resume_url if profile else None
        })
    
    return jsonify({
        "success": True,
        "drive": {
            "id": drive.id,
            "company_name": company.name if company else "Unknown",
            "title": drive.title,
            "description": drive.description,
            "role": drive.role,
            "ctc_min": drive.ctc_min,
            "ctc_max": drive.ctc_max,
            "location": drive.location,
            "drive_type": drive.drive_type,
            "eligibility_criteria": drive.eligibility_criteria,
            "application_start": drive.application_start.isoformat() if drive.application_start else None,
            "application_end": drive.application_end.isoformat() if drive.application_end else None,
            "status": drive.status,
            "created_at": drive.created_at.isoformat()
        },
        "applicants": applicants
    })

@app.route("/hrd/drives/<int:drive_id>/eligible-students", methods=["GET"])
@hrd_required
def get_eligible_students(drive_id):
    """Get eligible students for a drive based on criteria"""
    drive = PlacementDrive.query.get(drive_id)
    if not drive:
        return jsonify({"success": False, "message": "Drive not found"}), 404
    
    criteria = drive.eligibility_criteria
    cgpa_min = criteria.get("cgpa_min", 0)
    branches = criteria.get("branches", [])
    max_backlogs = criteria.get("max_backlogs", 10)
    year = criteria.get("year", 4)
    
    # Query students based on criteria
    query = User.query.filter_by(role="student", status="APPROVED")
    
    if branches:
        query = query.filter(User.degree.in_(branches))
    
    if year:
        query = query.filter_by(semester=year * 2)  # Assuming 2 semesters per year
    
    students = query.all()
    
    eligible = []
    for student in students:
        profile = StudentPlacementProfile.query.filter_by(student_id=student.id).first()
        
        eligible.append({
            "student_id": student.id,
            "name": student.name,
            "email": student.email,
            "srn": student.srn,
            "degree": student.degree,
            "semester": student.semester,
            "skills": profile.skills if profile else [],
            "resume_url": profile.resume_url if profile else None
        })
    
    return jsonify({"success": True, "eligible_students": eligible, "count": len(eligible)})

@app.route("/hrd/drives/<int:drive_id>", methods=["PUT"])
@hrd_required
def update_drive(drive_id):
    """Update drive details"""
    drive = PlacementDrive.query.get(drive_id)
    if not drive:
        return jsonify({"success": False, "message": "Drive not found"}), 404
    
    data = request.json or {}
    
    try:
        drive.title = data.get("title", drive.title)
        drive.description = data.get("description", drive.description)
        drive.role = data.get("role", drive.role)
        drive.ctc_min = data.get("ctc_min", drive.ctc_min)
        drive.ctc_max = data.get("ctc_max", drive.ctc_max)
        drive.location = data.get("location", drive.location)
        drive.status = data.get("status", drive.status)
        
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/hrd/applications/<int:app_id>/tag", methods=["PUT"])
@hrd_required
def tag_application(app_id):
    """Add tags to application"""
    application = DriveApplication.query.get(app_id)
    if not application:
        return jsonify({"success": False, "message": "Application not found"}), 404
    
    data = request.json or {}
    tags = data.get("tags", [])
    
    application.tags = tags
    db.session.commit()
    
    return jsonify({"success": True})

# ----- OFFER MANAGEMENT -----
@app.route("/hrd/offers", methods=["POST"])
@hrd_required
def create_offer():
    """Create placement offer"""
    data = request.json or {}
    hrd_id = get_jwt_identity()
    
    try:
        offer = PlacementOffer(
            drive_id=data.get("drive_id"),
            student_id=data.get("student_id"),
            role=data.get("role"),
            ctc=data.get("ctc"),
            location=data.get("location"),
            offer_letter_url=data.get("offer_letter_url"),
            joining_date=datetime.strptime(data.get("joining_date"), "%Y-%m-%d").date() if data.get("joining_date") else None,
            bond_details=data.get("bond_details"),
            expiry_date=datetime.strptime(data.get("expiry_date"), "%Y-%m-%d").date() if data.get("expiry_date") else None
        )
        db.session.add(offer)
        
        # Update application status
        application = DriveApplication.query.filter_by(
            drive_id=data.get("drive_id"),
            student_id=data.get("student_id")
        ).first()
        if application:
            application.status = "offer_received"
        
        db.session.commit()
        
        log_hrd_activity(hrd_id, "hrd", "created_offer", "offer", offer.id, {"student_id": offer.student_id, "ctc": offer.ctc})
        
        # Send email notification to student
        student = User.query.get(offer.student_id)
        if student:
            send_professional_email(
                student.email,
                "Placement Offer Received",
                "Congratulations on your Placement Offer!",
                {"Role": offer.role, "CTC": f"{offer.ctc} LPA", "Location": offer.location},
                f"You have received a placement offer for the role of {offer.role}."
            )
        
        return jsonify({"success": True, "offer_id": offer.id})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/hrd/offers", methods=["GET"])
@hrd_required
def get_offers():
    """Get all placement offers"""
    offers = PlacementOffer.query.all()
    
    result = []
    for offer in offers:
        student = User.query.get(offer.student_id)
        drive = PlacementDrive.query.get(offer.drive_id)
        company = Company.query.get(drive.company_id) if drive else None
        
        result.append({
            "id": offer.id,
            "student_name": student.name if student else "Unknown",
            "student_email": student.email if student else "",
            "company_name": company.name if company else "Unknown",
            "role": offer.role,
            "ctc": offer.ctc,
            "status": offer.status,
            "offer_date": offer.offer_date.isoformat(),
            "expiry_date": offer.expiry_date.isoformat() if offer.expiry_date else None
        })
    
    return jsonify({"success": True, "offers": result})

# ----- AI ANALYSIS -----
@app.route("/hrd/ai/analyze-resume", methods=["POST"])
@hrd_required
def analyze_resume_endpoint():
    """Analyze student resume using Groq AI"""
    data = request.json or {}
    resume_text = data.get("resume_text", "")
    student_id = data.get("student_id")
    
    if not resume_text:
        return jsonify({"success": False, "message": "Resume text required"}), 400
    
    # Extract skills
    skills_result = extract_skills_from_resume(resume_text)
    
    # Analyze quality
    quality_result = analyze_resume_quality(resume_text)
    
    # Cache results if student_id provided
    if student_id:
        cache = AIAnalysisCache.query.filter_by(student_id=student_id).first()
        if cache:
            cache.skills_extracted = skills_result.get("skills", [])
            cache.quality_score = quality_result.get("quality_score", 0)
            cache.ats_score = quality_result.get("ats_score", 0)
            cache.analyzed_at = datetime.utcnow()
        else:
            cache = AIAnalysisCache(
                student_id=student_id,
                resume_url=data.get("resume_url", ""),
                skills_extracted=skills_result.get("skills", []),
                quality_score=quality_result.get("quality_score", 0),
                ats_score=quality_result.get("ats_score", 0)
            )
            db.session.add(cache)
        db.session.commit()
    
    return jsonify({
        "success": True,
        "skills": skills_result.get("skills", []),
        "certifications": skills_result.get("certifications", []),
        "projects": skills_result.get("projects", []),
        "quality_score": quality_result.get("quality_score", 0),
        "ats_score": quality_result.get("ats_score", 0),
        "feedback": quality_result.get("feedback", "")
    })

@app.route("/hrd/ai/role-fit", methods=["POST"])
@hrd_required
def calculate_role_fit_endpoint():
    """Calculate role-fit for student"""
    data = request.json or {}
    student_skills = data.get("student_skills", [])
    job_requirements = data.get("job_requirements", [])
    
    result = calculate_role_fit(student_skills, job_requirements)
    return jsonify({"success": True, **result})

@app.route("/hrd/ai/batch-analysis", methods=["POST"])
@hrd_required
def batch_analysis():
    """Analyze skills across all students"""
    data = request.json or {}
    drive_id = data.get("drive_id")
    
    if drive_id:
        applications = DriveApplication.query.filter_by(drive_id=drive_id).all()
        student_ids = [app.student_id for app in applications]
    else:
        students = User.query.filter_by(role="student").all()
        student_ids = [s.id for s in students]
    
    all_skills = []
    for student_id in student_ids:
        profile = StudentPlacementProfile.query.filter_by(student_id=student_id).first()
        if profile and profile.skills:
            all_skills.append(profile.skills)
    
    analysis = analyze_batch_skills(all_skills)
    
    return jsonify({"success": True, **analysis})

# ----- ANALYTICS -----
@app.route("/hrd/analytics/overview", methods=["GET"])
@hrd_required
def analytics_overview():
    """Overall placement statistics"""
    total_students = User.query.filter_by(role="student", status="APPROVED").count()
    total_offers = PlacementOffer.query.count()
    accepted_offers = PlacementOffer.query.filter_by(status="accepted").count()
    pending_offers = PlacementOffer.query.filter_by(status="pending").count()
    
    active_drives = PlacementDrive.query.filter_by(status="open").count()
    
    # Calculate placement percentage
    placement_percentage = (accepted_offers / total_students * 100) if total_students > 0 else 0
    
    # Average CTC
    offers = PlacementOffer.query.filter_by(status="accepted").all()
    avg_ctc = sum([o.ctc for o in offers]) / len(offers) if offers else 0
    max_ctc = max([o.ctc for o in offers]) if offers else 0
    
    return jsonify({
        "success": True,
        "total_students": total_students,
        "total_offers": total_offers,
        "accepted_offers": accepted_offers,
        "pending_offers": pending_offers,
        "active_drives": active_drives,
        "placement_percentage": round(placement_percentage, 2),
        "avg_ctc": round(avg_ctc, 2),
        "max_ctc": max_ctc
    })

@app.route("/hrd/analytics/branch-wise", methods=["GET"])
@hrd_required
def analytics_branch_wise():
    """Branch-wise placement statistics"""
    degrees = db.session.query(User.degree).filter_by(role="student").distinct().all()
    
    result = []
    for (degree,) in degrees:
        total = User.query.filter_by(role="student", degree=degree).count()
        
        # Get placed students
        placed_student_ids = [o.student_id for o in PlacementOffer.query.filter_by(status="accepted").all()]
        placed_in_branch = User.query.filter(User.id.in_(placed_student_ids), User.degree == degree).count()
        
        placement_pct = (placed_in_branch / total * 100) if total > 0 else 0
        
        result.append({
            "branch": degree,
            "total": total,
            "placed": placed_in_branch,
            "percentage": round(placement_pct, 2)
        })
    
    return jsonify({"success": True, "branch_wise": result})

# ----- STUDENT PLACEMENT ROUTES -----
@app.route("/student/placement/profile", methods=["GET"])
@jwt_required()
def get_student_placement_profile():
    """Get student placement profile"""
    student_id = get_jwt_identity()
    
    profile = StudentPlacementProfile.query.filter_by(student_id=student_id).first()
    
    if not profile:
        return jsonify({"success": True, "profile": None})
    
    return jsonify({
        "success": True,
        "profile": {
            "resume_url": profile.resume_url,
            "skills": profile.skills,
            "certifications": profile.certifications,
            "projects": profile.projects,
            "github_url": profile.github_url,
            "linkedin_url": profile.linkedin_url,
            "portfolio_url": profile.portfolio_url,
            "preferred_location": profile.preferred_location,
            "preferred_role": profile.preferred_role,
            "min_expected_ctc": profile.min_expected_ctc,
            "profile_updated_at": profile.profile_updated_at.isoformat()
        }
    })

@app.route("/student/placement/profile", methods=["PUT"])
@jwt_required()
def update_student_placement_profile():
    """Update student placement profile"""
    student_id = get_jwt_identity()
    data = request.json or {}
    
    profile = StudentPlacementProfile.query.filter_by(student_id=student_id).first()
    
    try:
        if not profile:
            profile = StudentPlacementProfile(student_id=student_id)
            db.session.add(profile)
        
        profile.skills = data.get("skills", profile.skills if profile else [])
        profile.certifications = data.get("certifications", profile.certifications if profile else [])
        profile.projects = data.get("projects", profile.projects if profile else [])
        profile.github_url = data.get("github_url", profile.github_url if profile else None)
        profile.linkedin_url = data.get("linkedin_url", profile.linkedin_url if profile else None)
        profile.portfolio_url = data.get("portfolio_url", profile.portfolio_url if profile else None)
        profile.preferred_location = data.get("preferred_location", profile.preferred_location if profile else None)
        profile.preferred_role = data.get("preferred_role", profile.preferred_role if profile else None)
        profile.min_expected_ctc = data.get("min_expected_ctc", profile.min_expected_ctc if profile else None)
        profile.profile_updated_at = datetime.utcnow()
        
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/student/placement/drives", methods=["GET"])
@jwt_required()
def get_available_drives():
    """Get available placement drives for student"""
    student_id = get_jwt_identity()
    student = User.query.get(student_id)
    
    if not student:
        return jsonify({"success": False, "message": "Student not found"}), 404
    
    # Get all open drives
    drives = PlacementDrive.query.filter_by(status="open").all()
    
    result = []
    for drive in drives:
        company = Company.query.get(drive.company_id)
        
        # Check if already applied
        application = DriveApplication.query.filter_by(drive_id=drive.id, student_id=student_id).first()
        
        # Check eligibility
        criteria = drive.eligibility_criteria
        eligible = True
        
        if criteria.get("branches") and student.degree not in criteria.get("branches"):
            eligible = False
        
        result.append({
            "id": drive.id,
            "company_name": company.name if company else "Unknown",
            "title": drive.title,
            "description": drive.description,
            "role": drive.role,
            "ctc_min": drive.ctc_min,
            "ctc_max": drive.ctc_max,
            "location": drive.location,
            "application_end": drive.application_end.isoformat() if drive.application_end else None,
            "eligible": eligible,
            "applied": application is not None,
            "application_status": application.status if application else None
        })
    
    return jsonify({"success": True, "drives": result})

@app.route("/student/placement/drives/<int:drive_id>/apply", methods=["POST"])
@jwt_required()
def apply_to_drive(drive_id):
    """Apply to placement drive"""
    student_id = get_jwt_identity()
    
    drive = PlacementDrive.query.get(drive_id)
    if not drive:
        return jsonify({"success": False, "message": "Drive not found"}), 404
    
    if drive.status != "open":
        return jsonify({"success": False, "message": "Drive is not open for applications"}), 400
    
    # Check if already applied
    existing = DriveApplication.query.filter_by(drive_id=drive_id, student_id=student_id).first()
    if existing:
        return jsonify({"success": False, "message": "Already applied to this drive"}), 400
    
    try:
        application = DriveApplication(
            drive_id=drive_id,
            student_id=student_id
        )
        db.session.add(application)
        db.session.commit()
        
        log_hrd_activity(student_id, "student", "applied_to_drive", "application", application.id, {"drive_id": drive_id})
        
        return jsonify({"success": True, "application_id": application.id})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/student/placement/offers", methods=["GET"])
@jwt_required()
def get_student_offers():
    """Get student's placement offers"""
    student_id = get_jwt_identity()
    
    offers = PlacementOffer.query.filter_by(student_id=student_id).all()
    
    result = []
    for offer in offers:
        drive = PlacementDrive.query.get(offer.drive_id)
        company = Company.query.get(drive.company_id) if drive else None
        
        result.append({
            "id": offer.id,
            "company_name": company.name if company else "Unknown",
            "role": offer.role,
            "ctc": offer.ctc,
            "location": offer.location,
            "offer_letter_url": offer.offer_letter_url,
            "joining_date": offer.joining_date.isoformat() if offer.joining_date else None,
            "bond_details": offer.bond_details,
            "offer_date": offer.offer_date.isoformat(),
            "expiry_date": offer.expiry_date.isoformat() if offer.expiry_date else None,
            "status": offer.status
        })
    
    return jsonify({"success": True, "offers": result})

@app.route("/student/placement/offers/<int:offer_id>/accept", methods=["POST"])
@jwt_required()
def accept_offer(offer_id):
    """Accept placement offer"""
    student_id = get_jwt_identity()
    
    offer = PlacementOffer.query.get(offer_id)
    if not offer or offer.student_id != student_id:
        return jsonify({"success": False, "message": "Offer not found"}), 404
    
    if offer.status != "pending":
        return jsonify({"success": False, "message": "Offer already responded to"}), 400
    
    try:
        offer.status = "accepted"
        offer.student_response_date = datetime.utcnow()
        
        # Update application status
        application = DriveApplication.query.filter_by(
            drive_id=offer.drive_id,
            student_id=student_id
        ).first()
        if application:
            application.status = "offer_accepted"
        
        db.session.commit()
        
        log_hrd_activity(student_id, "student", "accepted_offer", "offer", offer_id, {"ctc": offer.ctc})
        
        return jsonify({"success": True})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/student/placement/offers/<int:offer_id>/reject", methods=["POST"])
@jwt_required()
def reject_offer(offer_id):
    """Reject placement offer"""
    student_id = get_jwt_identity()
    data = request.json or {}
    
    offer = PlacementOffer.query.get(offer_id)
    if not offer or offer.student_id != student_id:
        return jsonify({"success": False, "message": "Offer not found"}), 404
    
    if offer.status != "pending":
        return jsonify({"success": False, "message": "Offer already responded to"}), 400
    
    try:
        offer.status = "rejected"
        offer.student_response_date = datetime.utcnow()
        offer.rejection_reason = data.get("reason", "")
        
        # Update application status
        application = DriveApplication.query.filter_by(
            drive_id=offer.drive_id,
            student_id=student_id
        ).first()
        if application:
            application.status = "offer_rejected"
        
        db.session.commit()
        
        log_hrd_activity(student_id, "student", "rejected_offer", "offer", offer_id, {"reason": offer.rejection_reason})
        
        return jsonify({"success": True})
    except Exception as e:
        db.session.rollback()
        return jsonify({"success": False, "message": str(e)}), 500

# Register Blueprint
app.register_blueprint(attendance_bp, url_prefix='/attendance')

if __name__ == "__main__":
    with app.app_context():
        init_db()
        print("Using GROQ_API_KEY:", (GROQ_API_KEY[:8] + "********") if GROQ_API_KEY else "NOT SET")
        print("JWT_SECRET_KEY loaded:", True if JWT_SECRET_KEY else False)
    app.run(debug=True, host='0.0.0.0', port=FLASK_RUN_PORT)
